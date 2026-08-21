"""Microsoft-style theme helpers for the Agent Insights deck."""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------- palette
NAVY = RGBColor(0x14, 0x24, 0x3E)
NAVY_2 = RGBColor(0x1E, 0x33, 0x54)
BLUE = RGBColor(0x00, 0x78, 0xD4)
BLUE_L = RGBColor(0x50, 0xE6, 0xFF)
BLUE_XL = RGBColor(0xDE, 0xEC, 0xF9)
PURPLE = RGBColor(0x86, 0x61, 0xC5)
TEAL = RGBColor(0x00, 0xB7, 0xC3)
MAGENTA = RGBColor(0xE3, 0x00, 0x8C)
AMBER = RGBColor(0xE8, 0x7C, 0x00)
GREEN = RGBColor(0x0E, 0x70, 0x0E)
RED = RGBColor(0xD1, 0x34, 0x38)
INK = RGBColor(0x1B, 0x1A, 0x19)
GRAY = RGBColor(0x60, 0x5E, 0x5C)
GRAY_L = RGBColor(0x8A, 0x88, 0x86)
LINE = RGBColor(0xE1, 0xDF, 0xDD)
WASH = RGBColor(0xF7, 0xF7, 0xF9)
WASH_2 = RGBColor(0xF0, 0xF3, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Segoe UI"
FONT_LT = "Segoe UI Light"
FONT_SB = "Segoe UI Semibold"
MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.72)
CONTENT_W = Inches(11.89)


# ---------------------------------------------------------------- helpers
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def rect(slide, x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE,
         radius=None):
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    if radius is not None:
        try:
            sh.adjustments[0] = radius
        except (IndexError, KeyError):
            pass
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    sh.text_frame.word_wrap = True
    return sh


def card(slide, x, y, w, h, fill=WHITE, line=LINE, radius=0.045):
    return rect(slide, x, y, w, h, fill=fill, line=line, lw=1.0,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=radius)


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf


def para(tf, text, size=14, color=INK, bold=False, font=FONT, first=False,
         space_before=0, space_after=4, align=PP_ALIGN.LEFT, line=None, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if line:
        p.line_spacing = line
    if isinstance(text, str):
        text = [(text, {})]
    for chunk, opt in text:
        r = p.add_run()
        r.text = chunk
        f = r.font
        f.name = opt.get("font", font)
        f.size = Pt(opt.get("size", size))
        f.bold = opt.get("bold", bold)
        f.italic = opt.get("italic", italic)
        f.color.rgb = opt.get("color", color)
    return p


def fill_shape_text(sh, lines, anchor=MSO_ANCHOR.MIDDLE, pad=Inches(0.1)):
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = pad
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    for i, (text, kw) in enumerate(lines):
        para(tf, text, first=(i == 0), **kw)
    return sh


def bg(slide, color):
    sh = rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=color)
    slide.shapes._spTree.remove(sh._element)
    slide.shapes._spTree.insert(2, sh._element)
    return sh


def accent_bar(slide, y=Inches(0.0), colors=(BLUE, TEAL, PURPLE, MAGENTA), h=Inches(0.075)):
    w = SLIDE_W / len(colors)
    for i, c in enumerate(colors):
        rect(slide, Emu(int(i * w)), y, Emu(int(w) + 2), h, fill=c)


def header(slide, kicker, title, sub=None, rule=True):
    """Standard content-slide header. Returns the y where body content starts."""
    accent_bar(slide)
    tf = textbox(slide, MARGIN, Inches(0.42), CONTENT_W, Inches(0.26))
    para(tf, kicker.upper(), size=10.5, color=BLUE, bold=True, font=FONT_SB, first=True,
         space_after=0)
    tf = textbox(slide, MARGIN, Inches(0.72), CONTENT_W, Inches(0.55))
    para(tf, title, size=29, color=NAVY, font=FONT_LT, first=True, space_after=0)
    y = Inches(1.42)
    if sub:
        tf = textbox(slide, MARGIN, Inches(1.32), CONTENT_W, Inches(0.36))
        para(tf, sub, size=13.5, color=GRAY, first=True, space_after=0, line=1.2)
        y = Inches(1.9)
    if rule:
        rect(slide, MARGIN, y - Inches(0.16), Inches(1.05), Inches(0.035), fill=BLUE)
    return y


PAGE_TOKEN = "{{N}}"
FOOT = "Microsoft  \u00b7  Agent Insights"


def footer(slide, text, num=None):
    """Draw the footer. The page number is a token resolved at save time by
    number_pages(), so inserting or reordering slides never desynchronises it."""
    tf = textbox(slide, MARGIN, Inches(6.98), Inches(9.5), Inches(0.26))
    para(tf, text, size=9, color=GRAY_L, first=True, space_after=0)
    tf = textbox(slide, Inches(11.85), Inches(6.98), Inches(0.76), Inches(0.26))
    para(tf, PAGE_TOKEN, size=9, color=GRAY_L, first=True, space_after=0,
         align=PP_ALIGN.RIGHT)


def number_pages(prs):
    """Replace every PAGE_TOKEN with the slide's real 1-based position."""
    for i, slide in enumerate(prs.slides, start=1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.text == PAGE_TOKEN:
                        r.text = str(i)


def down_arrow(slide, cx, y, color=RGBColor(0xC8, 0xC6, 0xC4), w=Inches(0.2), h=Inches(0.2)):
    rect(slide, Emu(int(cx - w / 2)), y, w, h, fill=color, shape=MSO_SHAPE.DOWN_ARROW)


def chip(slide, x, y, w, h, label, fill, fg=WHITE, size=11, bold=True):
    sh = card(slide, x, y, w, h, fill=fill, line=fill, radius=0.5)
    fill_shape_text(sh, [(label, dict(size=size, color=fg, bold=bold,
                                      align=PP_ALIGN.CENTER, space_after=0))])
    return sh


def style_table(tbl, header_fill=NAVY, header_fg=WHITE, size=10.5, header_size=10.5,
                row_fills=(WHITE, WASH), first_col_bold=True, first_col_color=NAVY):
    for r, row in enumerate(tbl.rows):
        for c, cell in enumerate(row.cells):
            cell.margin_left = Inches(0.09)
            cell.margin_right = Inches(0.09)
            cell.margin_top = Inches(0.045)
            cell.margin_bottom = Inches(0.045)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if r == 0 else row_fills[(r - 1) % len(row_fills)]
            for p in cell.text_frame.paragraphs:
                p.space_after = Pt(0)
                p.space_before = Pt(0)
                p.line_spacing = 1.05
                for run in p.runs:
                    run.font.name = FONT
                    run.font.size = Pt(header_size if r == 0 else size)
                    if r == 0:
                        run.font.bold = True
                        run.font.color.rgb = header_fg
                    elif c == 0 and first_col_bold:
                        run.font.bold = True
                        run.font.color.rgb = first_col_color
                    else:
                        run.font.color.rgb = INK
    return tbl


def add_table(slide, rows, x, y, w, col_w=None, row_h=Inches(0.34), header_h=Inches(0.38)):
    n_r, n_c = len(rows), len(rows[0])
    gt = slide.shapes.add_table(n_r, n_c, x, y, w, header_h + row_h * (n_r - 1))
    tbl = gt.table
    tbl.first_row = True
    tbl.horz_banding = False
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Emu(int(w * cw / total))
    tbl.rows[0].height = header_h
    for i in range(1, n_r):
        tbl.rows[i].height = row_h
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            tbl.cell(r, c).text = val
    return tbl


def code_block(slide, x, y, w, h, lines, size=11.5, fill=NAVY, fg=RGBColor(0xE8, 0xEE, 0xF7)):
    sh = card(slide, x, y, w, h, fill=fill, line=fill, radius=0.03)
    tf = sh.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = Inches(0.22)
    tf.margin_top = tf.margin_bottom = Inches(0.16)
    for i, ln in enumerate(lines):
        if isinstance(ln, str):
            ln = [(ln, {})]
        para(tf, ln, size=size, color=fg, font=MONO, first=(i == 0), space_after=1,
             line=1.16)
    return sh
