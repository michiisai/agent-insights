"""New slides adopting structure from the Explore 2025 intern deck."""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from theme import *



def photo_slot(slide, cx, cy, d, label="add photo"):
    """Dashed circular placeholder the presenter can fill with a headshot."""
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(cx - d / 2)), Emu(int(cy - d / 2)), d, d)
    sh.fill.solid()
    sh.fill.fore_color.rgb = WASH_2
    sh.line.color.rgb = RGBColor(0xB3, 0xC4, 0xDC)
    sh.line.width = Pt(1.25)
    sh.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, label, size=9, color=GRAY_L, align=PP_ALIGN.CENTER, first=True, space_after=0)
    return sh


# ------------------------------------------------------------------ about me
def s_about(prs):
    s = blank(prs)
    y = header(s, "Introduction", "About me")

    photo_slot(s, Inches(2.35), Inches(3.30), Inches(2.1), "add\nheadshot")

    tf = textbox(s, Inches(3.72), Inches(1.92), Inches(4.3), Inches(1.0))
    para(tf, "Michelle Ma", size=32, color=NAVY, font=FONT_LT, first=True, space_after=4)
    para(tf, "Software Engineer Intern  \u00b7  VS Code", size=13.5, color=BLUE, bold=True,
         space_after=0)

    rows = [("\U0001F393", "UC San Diego", "Mathematics, Computer Science & Cognitive Science", BLUE),
            ("\U0001F4C5", "2x Microsoft Intern", "SWE + Explore", PURPLE),
            ("\U0001F50D", "Favorite Memory", "Bainbridge Island", TEAL),
            ("\u2B50", "Fun Fact", "I have journaled every day for the past 600+ days", MAGENTA)]
    yy = Inches(3.10)
    for icon, title, sub, col in rows:
        c = card(s, Inches(3.72), yy, Inches(8.89), Inches(0.72), fill=WHITE, line=LINE)
        rect(s, Inches(3.72), yy, Inches(0.06), Inches(0.72), fill=col)
        tf = textbox(s, Inches(4.0), yy + Inches(0.19), Inches(0.4), Inches(0.4))
        para(tf, icon, size=14, first=True, space_after=0)
        tf = textbox(s, Inches(4.48), yy + Inches(0.12), Inches(7.9), Inches(0.55))
        para(tf, title, size=12.5, color=NAVY, bold=True, first=True, space_after=2)
        para(tf, sub, size=11, color=GRAY, space_after=0, line=1.16)
        yy += Inches(0.82)

    footer(s, FOOT)
    notes(s, """
SCRIPT

Twenty seconds. This is a handshake, not a biography. Do not linger.

"Quickly about me \u2014 I'm at UC San Diego studying math, computer science and cognitive science,
and this is my second internship here. Explore last year, and software engineering on VS Code
this year."

[pick ONE of the bottom two rows \u2014 don't read both]

"And a fun fact: I've journaled every day for the past six hundred days."

[move on immediately]

\u2500\u2500\u2500 NOTES \u2500\u2500\u2500

\u2022 The journaling line is the better one if the room is quiet and you want a beat of personality.
  Bainbridge Island is the better one if you want to connect with people who were on that trip.
\u2022 Resist explaining why you picked this project here \u2014 that's what the next four slides are for,
  and the answer lands much harder after they've felt the problem.
\u2022 CHECK BEFORE PRESENTING: swap in your own headshot.
""")
    return s


# ------------------------------------------------------------------ team / relevance
def s_context(prs):
    s = blank(prs)
    y = header(s, "01  Problem \u00b7 where it lives", "Why this is a VS Code problem",
               "The team ships the editor most of these agents now run inside.")

    sh = card(s, MARGIN, y, Inches(11.89), Inches(1.06), fill=NAVY, line=NAVY)
    rect(s, MARGIN, y, Inches(0.075), Inches(1.06), fill=TEAL)
    tf = sh.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.42)
    tf.margin_right = Inches(0.36)
    para(tf, [("VS Code is where coding agents actually run. ", dict(color=WHITE, bold=True)),
              ("That makes the editor the only place on the machine that can see all of them "
               "at once \u2014 and the only place a developer would trust with what they say.",
               dict(color=RGBColor(0xB9, 0xC8, 0xE0)))],
         size=16, first=True, space_after=0, line=1.25)

    cards_ = [
        ("The editor became the host", BLUE,
         "Agent Host runs Claude Code and Codex as sessions inside VS Code, alongside Copilot."),
        ("VS Code is already in the path", TEAL,
         "The host routes their telemetry and stamps its own conversation id onto it."),
        ("Agent cost became a real concern", AMBER,
         "Tokens, cache hit rate and failed runs are day-to-day questions now \u2014 "
         "including for this team."),
        ("It has to stay local", PURPLE,
         "Prompts, model replies, tool arguments, source code. Exactly what you cannot send "
         "to a shared collector."),
    ]
    cw = Inches(2.86)
    gap = Inches(0.155)
    yy = y + Inches(1.32)
    for i, (title, col, body) in enumerate(cards_):
        x = MARGIN + (cw + gap) * i
        c = card(s, x, yy, cw, Inches(2.1), fill=WHITE, line=LINE)
        rect(s, x, yy, cw, Inches(0.055), fill=col)
        tf = c.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = tf.margin_right = Inches(0.24)
        tf.margin_top = Inches(0.24)
        para(tf, title, size=12.5, color=NAVY, bold=True, first=True, space_after=7, line=1.14)
        para(tf, body, size=10.5, color=GRAY, space_after=0, line=1.22)

    band = card(s, MARGIN, Inches(5.86), Inches(11.89), Inches(0.8), fill=WASH_2, line=WASH_2)
    tf = band.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.34)
    para(tf, [("So the observability has to live in the editor too. ", dict(bold=True, color=BLUE)),
              ("Not in a cloud dashboard, not in a separate app \u2014 next to the thing it describes, "
               "on your own machine.", dict(color=NAVY))],
         size=13, first=True, space_after=0, line=1.22)

    footer(s, FOOT)
    notes(s, """
SCRIPT

Thirty seconds. This slide opens the Problem section by establishing WHERE the problem lives.
Don't rush it, but don't add to it either.

"I spent the summer on VS Code \u2014 the team that builds the editor, in the open, shipping monthly."

[the navy band \u2014 this is the actual argument, say it close to verbatim]
"And VS Code is where coding agents actually run now. That makes the editor the only place on
the machine that can see all of them at once \u2014 and the only place a developer would trust with
what they say."

[the four cards \u2014 ONE clause each, do not read the body text]
"The editor became the host \u2014 Claude Code and Codex run as sessions inside VS Code now, next to
Copilot."

"And VS Code isn't just next to that telemetry, it's already in the path. It routes it, and
stamps its own conversation id onto it."

"Agent cost stopped being an abstraction \u2014 it's a question this team asks about its own work."

"But the contents are prompts, model replies and source code. So it can't leave the machine."

[the bottom band]
"Which means the observability has to live in the editor too. Not in a cloud dashboard, not in a
separate app \u2014 next to the thing it describes."

\u2500\u2500\u2500 NOTES \u2500\u2500\u2500

\u2022 Card 2 is the hand-off. You've just said telemetry is flowing through the host, so the next
  slide spends thirty seconds on what that telemetry actually is. Don't explain it here.
\u2022 If you want a number, VS Code's install base is a fine one to say out loud \u2014 but keep it off
  the slide. The argument is stronger than the statistic.
\u2022 If asked "why not a cloud service?" \u2014 that's the fourth card, and the appendix has a full
  data-flow diagram showing the trust boundary.
""")
    return s


def s_goal(prs):
    s = blank(prs)
    y = header(s, "01  Problem \u00b7 the brief", "What I set out to build",
               "One extension that makes agent telemetry readable \u2014 without sending "
               "any of it off your machine.")

    # the goal statement
    sh = card(s, MARGIN, y, Inches(11.89), Inches(1.16), fill=NAVY, line=NAVY)
    rect(s, MARGIN, y, Inches(0.075), Inches(1.16), fill=BLUE)
    tf = sh.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.42)
    tf.margin_right = Inches(0.36)
    para(tf, "THE GOAL", size=10, color=BLUE_L, bold=True, font=FONT_SB, first=True, space_after=7)
    para(tf, [("Close the gap ", dict(color=WHITE)),
              ("after", dict(color=BLUE_L, bold=True, italic=True)),
              (" you turn OpenTelemetry on \u2014 give humans and agents inside VS Code a way to "
               "actually consume the telemetry their agent runs already emit.", dict(color=WHITE))],
         size=16, space_after=0, line=1.25)

    # what the extension had to do
    pieces = [
        ("Collect", BLUE, "Receive OTLP from every agent on the machine \u2014 without any of it "
                          "leaving the machine"),
        ("Store", TEAL, "Keep it locally, queryable, and small enough to live inside an editor"),
        ("Show a person", AMBER, "Explorers for traces, logs and metrics \u2014 and a view of an "
                                 "agent session as a session"),
        ("Answer an agent", PURPLE, "A tool surface Copilot can call, so the answer comes back "
                                    "in chat with a link into the panel"),
    ]
    cw = Inches(2.86)
    gap = Inches(0.155)
    yy = y + Inches(1.42)
    for i, (title, col, body) in enumerate(pieces):
        x = MARGIN + (cw + gap) * i
        c = card(s, x, yy, cw, Inches(1.94), fill=WHITE, line=LINE)
        rect(s, x, yy, cw, Inches(0.055), fill=col)
        nb = card(s, x + Inches(0.24), yy + Inches(0.3), Inches(0.36), Inches(0.36),
                  fill=col, line=col, radius=0.5)
        fill_shape_text(nb, [(str(i + 1), dict(size=11.5, color=WHITE, bold=True,
                                               align=PP_ALIGN.CENTER, space_after=0))])
        tf = textbox(s, x + Inches(0.7), yy + Inches(0.31), cw - Inches(0.94), Inches(0.36))
        para(tf, title, size=13.5, color=NAVY, bold=True, first=True, space_after=0)
        tf = textbox(s, x + Inches(0.24), yy + Inches(0.82), cw - Inches(0.48), Inches(1.0))
        para(tf, body, size=10.5, color=GRAY, first=True, space_after=0, line=1.22)

    band = card(s, MARGIN, Inches(5.86), Inches(11.89), Inches(0.8), fill=BLUE_XL, line=BLUE_XL)
    tf = band.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.34)
    para(tf, [("Four verbs, one extension, one machine. ", dict(bold=True, color=BLUE)),
              ("The hard part turned out to be a fifth step nobody puts on the list \u2014 "
               "deciding what the data actually means.",
               dict(color=NAVY))],
         size=13, first=True, space_after=0, line=1.22)

    footer(s, FOOT)
    notes(s, """
SCRIPT

Twenty-five seconds. This slide is a signpost, not a destination. Keep it brisk.

"That's the problem. Here's what I set out to build."

[read the navy statement once, slowly \u2014 it's the thesis in project form]
"Close the gap after you turn OpenTelemetry on. Give humans and agents inside VS Code a way to
actually consume the telemetry their agent runs already emit."

[the four cards \u2014 four verbs, one breath total, do NOT read the body text]
"Collect it. Store it. Show it to a person. Answer questions for an agent. One extension,
everything local."

[the band \u2014 land it and MOVE ON]
"Four verbs. And the hard part turned out to be a fifth step nobody puts on the list \u2014 deciding
what the data actually means."

\u2500\u2500\u2500 NOTES \u2500\u2500\u2500

\u2022 That last clause is the setup for the entire Approach section. You pay it off two slides later
  on "seven things have to happen". Don't explain it here.
\u2022 If asked about VS Code emitting its own cross-process telemetry: that's on the roadmap slide.
  The confident answer is "it's the natural next surface, and this extension is exactly the
  consumer it needs, so it's shovel-ready rather than hypothetical." Don't pre-empt it here.
""")
    return s


# ------------------------------------------------------------------ pipeline
PIPELINE = [
    ("1", "Emit", "The harness writes spans, logs and metrics over OTLP"),
    ("2", "Group into traces", "Spans that share a trace id become one trace"),
    ("3", "Identify the session", "Which conversation does this trace belong to?"),
    ("4", "Split into turns", "One user request and the model calls it triggered"),
    ("5", "Count the work", "How many model calls, how many tool calls"),
    ("6", "Reconcile tokens", "Input, output, cache read, cache creation"),
    ("7", "Rebuild the transcript", "What the person and the model actually said"),
]

FRICTION = [
    "Codex strips assistant text; Claude splits the two halves of a conversation across two signals",
    "The host anchors only the trace that started the thread",
    "Codex carries no id at all; MAX(service_name) relabels every session \u201cvscode-agent-host\u201d",
    "Claude logs api_request when the round trip finishes \u2014 records arrive before the call they close",
    "One Codex call = 5 nested spans; matching Claude\u2019s subtree trebles the tool count",
    "Three incompatible conventions, plus parent spans that re-report their children",
    "Three completely different sources, one of which never sends the model\u2019s half",
]


def _pipeline(s, y, friction=False):
    bw = Inches(2.8)
    gap = Inches(0.23)
    bh = Inches(1.34) if friction else Inches(1.06)
    for i, (num, title, sub) in enumerate(PIPELINE):
        row, col_i = divmod(i, 4)
        x = MARGIN + (bw + gap) * col_i
        yy = y + (bh + Inches(0.46)) * row
        accent = RED if friction else BLUE
        c = card(s, x, yy, bw, bh, fill=WHITE, line=LINE)
        rect(s, x, yy, bw, Inches(0.05), fill=accent)
        nb = card(s, x + Inches(0.18), yy + Inches(0.2), Inches(0.34), Inches(0.34),
                  fill=accent, line=accent, radius=0.5)
        fill_shape_text(nb, [(num, dict(size=11, color=WHITE, bold=True,
                                        align=PP_ALIGN.CENTER, space_after=0))])
        tf = textbox(s, x + Inches(0.62), yy + Inches(0.2), bw - Inches(0.82), Inches(0.36))
        para(tf, title, size=12.5, color=NAVY, bold=True, first=True, space_after=0)
        tf = textbox(s, x + Inches(0.18), yy + Inches(0.62), bw - Inches(0.36),
                     bh - Inches(0.72))
        if friction:
            para(tf, FRICTION[i], size=9.5, color=RED, first=True, space_after=0, line=1.18)
        else:
            para(tf, sub, size=10.5, color=GRAY, first=True, space_after=0, line=1.18)
        if col_i < 3 and i < len(PIPELINE) - 1:
            rect(s, x + bw + Inches(0.045), yy + Inches(0.3), Inches(0.14), Inches(0.18),
                 fill=RGBColor(0xC8, 0xC6, 0xC4), shape=MSO_SHAPE.RIGHT_ARROW)

    # result chip in the empty 8th cell
    x = MARGIN + (bw + gap) * 3
    yy = y + (bh + Inches(0.46))
    fill = NAVY if not friction else RGBColor(0x6B, 0x1A, 0x1E)
    c = card(s, x, yy, bw, bh, fill=fill, line=fill)
    tf = c.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.2)
    if friction:
        para(tf, "Every stage breaks", size=14, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, first=True, space_after=5)
        para(tf, "\u2026and it breaks differently for each of the three harnesses",
             size=10.5, color=RGBColor(0xF2, 0xC5, 0xC8), align=PP_ALIGN.CENTER,
             space_after=0, line=1.2)
    else:
        para(tf, "A readable session", size=14, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, first=True, space_after=5)
        para(tf, "outcome \u00b7 timeline \u00b7 tokens \u00b7 transcript", size=10.5, color=BLUE_L,
             align=PP_ALIGN.CENTER, space_after=0, line=1.2)


def s_pipeline(prs):
    s = blank(prs)
    y = header(s, "03  Approach \u00b7 the shape of the work",
               "Seven things have to happen before you can read a session",
               "None of these are given to you. Every one is something the engine has to decide.")
    _pipeline(s, y + Inches(0.16), friction=False)
    band = card(s, MARGIN, Inches(5.25), CONTENT_W, Inches(0.86),
                fill=WASH_2, line=WASH_2)
    tf = band.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.3)
    para(tf, [("Every box is a decision, not a lookup.  ", dict(color=BLUE, bold=True)),
              ("The OpenTelemetry GenAI conventions describe how to write this data. "
               "Nothing describes how to read it back as a conversation.", {})],
         size=13.5, color=NAVY, first=True, space_after=0, line=1.25)
    footer(s, FOOT)
    notes(s, """
SCRIPT

Twenty-five seconds. Establish the map \u2014 the NEXT slide is the payoff, so don't spend here.

"Before any of this is a product, seven things have to happen to turn raw telemetry into
something you'd actually read."

[walk 1 \u2192 7 with your hand, NAMING them only \u2014 don't explain any of them yet]
"Emit it. Group it into traces. Work out which session a trace belongs to. Split that into
turns. Count the work. Reconcile the tokens. Rebuild the transcript."

[the navy box]
"And out the other end you get a readable session \u2014 an outcome, a timeline, tokens, and a
transcript."

[the band]
"And every one of those boxes is a decision, not a lookup. The OpenTelemetry conventions
describe how to write this data. Nothing describes how to read it back as a conversation."

Then set up the next slide and advance immediately:
"That's the happy path. Here's what actually happens."

\u2500\u2500\u2500 NOTES \u2500\u2500\u2500

\u2022 Resist explaining any individual stage here. The whole value of this slide is that the next
  one re-shows it in red \u2014 and that only works if the audience has the map first.
\u2022 If you're behind time, this pair is still worth keeping. It's the cheapest way to make the
  normalization argument land.
""")
    return s


def s_friction(prs):
    s = blank(prs)
    y = header(s, "03  Approach \u00b7 the shape of the work",
               "\u2026and every one of them breaks",
               "Same seven stages. This is where the three harnesses disagree with each other.")
    _pipeline(s, y + Inches(0.16), friction=True)
    band = card(s, MARGIN, Inches(5.55), CONTENT_W, Inches(0.86),
                fill=RGBColor(0xFD, 0xF3, 0xF4), line=RGBColor(0xF2, 0xD8, 0xDA))
    tf = band.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.3)
    para(tf, [("Roughly 2,600 lines of the engine exist for this slide.  ",
               dict(color=RED, bold=True)),
              ("The receiver that accepts the telemetry in the first place is about 200.", {})],
         size=13.5, color=NAVY, first=True, space_after=0, line=1.25)
    footer(s, FOOT)
    notes(s, """
SCRIPT

Twenty-five seconds, and it costs almost nothing \u2014 they already know the map.

The whole point of this slide is the RE-SHOW. Same seven boxes, now all red. Let the visual do
the work. Do NOT read all seven.

"Same seven stages. And every single one of them breaks \u2014 differently, depending on which agent
produced the data."

[pick exactly TWO and say them out loud]

"Stage three, identify the session. Codex frequently carries no conversation id at all. And the
agent host relabels every session 'vscode-agent-host' \u2014 for no better reason than that it sorts
last alphabetically and I was taking a maximum."

"Stage six, tokens. Three incompatible counting conventions co-existing on one machine, plus
parent spans that re-report everything their children already reported."

[the band \u2014 this is the number that reframes the project]
"Roughly two thousand six hundred lines of the engine exist for this slide. The receiver that
accepts the telemetry in the first place is about two hundred."

[close, then advance into the architecture]
"Everything after this is how I made those seven stages produce one answer instead of three."

\u2500\u2500\u2500 NOTES \u2500\u2500\u2500

\u2022 The 200-vs-2,600 ratio is the single most memorable fact in the deck. It comes back on the
  learning outcomes slide \u2014 planting it here makes that callback work.
\u2022 If someone asks about a stage you didn't cover, every one of them has a real story. Stage 1,
  Codex strips assistant text. Stage 2, the host anchors only the trace that started the thread.
  Stage 5, one Codex call is five spans. Stage 7, three different sources for the transcript.
""")
    return s


# ------------------------------------------------------------------ real use cases
def s_use_cases(prs):
    s = blank(prs)
    y = header(s, "02  The twist \u00b7 not just me", "This is a live need, not a hypothesis",
               "Three independent signals that people want to see what their agents are doing.")

    # ---- left: evidence screenshot
    lw = Inches(6.5)
    c = card(s, MARGIN, y, lw, Inches(3.62), fill=WHITE, line=LINE)
    rect(s, MARGIN, y, lw, Inches(0.055), fill=BLUE)
    tf = textbox(s, MARGIN + Inches(0.28), y + Inches(0.24), lw - Inches(0.56), Inches(0.34))
    para(tf, "People are already filing this as a bug", size=12.5, color=NAVY, bold=True,
         first=True, space_after=0)
    ph = rect(s, MARGIN + Inches(0.28), y + Inches(0.72), lw - Inches(0.56), Inches(2.7),
              fill=WASH_2, line=RGBColor(0xB3, 0xC4, 0xDC), lw=1.25)
    ph.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    tf = ph.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, "\U0001F5BC", size=22, color=RGBColor(0xB3, 0xC4, 0xDC),
         align=PP_ALIGN.CENTER, first=True, space_after=6)
    para(tf, "Screenshot the Agent Debug Logs issue \u2014 crop to the title and\n"
             "\u201cActual: no Claude agent harness session logs found\u201d",
         size=11, color=GRAY_L, align=PP_ALIGN.CENTER, space_after=0, line=1.24)

    # ---- right: three signals
    rx = MARGIN + lw + Inches(0.23)
    rw = Inches(5.16)
    signals = [
        ("The in-box view has gaps", BLUE,
         "Agent-host sessions don\u2019t reliably appear in Agent Debug Logs \u2014 people are hitting "
         "it and reporting it."),
        ("Three vendors instrumented independently", PURPLE,
         "Copilot, Claude Code and Codex each decided on their own that agent runs are worth "
         "emitting OpenTelemetry for."),
        ("And the host is still rolling out", TEAL,
         "More agents running inside the editor, on more machines, every week \u2014 with the same "
         "visibility gap."),
    ]
    sy = y + Inches(0.02)
    for title, col, body in signals:
        cc = card(s, rx, sy, rw, Inches(1.14), fill=WHITE, line=LINE)
        rect(s, rx, sy, Inches(0.06), Inches(1.14), fill=col)
        tf = cc.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.24)
        para(tf, title, size=12.5, color=NAVY, bold=True, first=True, space_after=4, line=1.14)
        para(tf, body, size=10.5, color=GRAY, space_after=0, line=1.2)
        sy += Inches(1.26)

    band = card(s, MARGIN, Inches(5.72), Inches(11.89), Inches(0.94), fill=NAVY, line=NAVY)
    rect(s, MARGIN, Inches(5.72), Inches(0.075), Inches(0.94), fill=GREEN)
    tf = band.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.4)
    para(tf, [("Nobody has to take my word for the gap. ", dict(color=BLUE_L, bold=True)),
              ("And the question people ask is never \u201cshow me a span\u201d \u2014 it is "
               "\u201cwhat did it do?\u201d, which is exactly the question the session model exists to answer.",
               dict(color=RGBColor(0xC3, 0xD1, 0xE6)))],
         size=13, first=True, space_after=0, line=1.24)

    footer(s, FOOT)
    notes(s, """
SCRIPT

Thirty seconds. You've just spent a slide on what the host broke FOR YOU. This one shows it
isn't only you \u2014 which is what turns a personal war story into a product argument.

[open by tying straight back to the previous slide]
"That second problem \u2014 the in-box view not keeping up \u2014 isn't something only I noticed."

[the screenshot]
"This is a bug someone on the team filed. Agent Debug Logs don't show Claude harness sessions at
all. Expected: session logs. Actual: nothing."

[the three signals, one sentence each]
"So \u2014 people are filing bugs because they can't see agent-host sessions in the in-box view."

"All three vendors independently decided this data was worth emitting. That's three separate
teams agreeing the problem exists."

"And the host is still rolling out, so there's more of this every week, not less."

[the band]
"Nobody has to take my word for the gap. And the question people ask is never 'show me a span' \u2014
it's 'what did it do?'. Which is exactly the question the session model exists to answer."

[hand off into Approach]
"So that's the problem, and it got bigger while I was working on it. Here's what I built."

\u2500\u2500\u2500 NOTES \u2500\u2500\u2500

\u2022 Be generous about the debug logs. It's a gap in something shipping fast, not a failing, and it
  will get fixed. Do not dunk.
\u2022 This is deliberately NOT a "look how many people used my extension" slide. It was a summer
  project with limited dogfooding, and claiming adoption you don't have is the fastest way to
  lose a room. What you CAN claim honestly is that the need is real and independently evidenced.
\u2022 TO DO BEFORE PRESENTING: screenshot the issue into the left frame. Crop to the title and the
  "Actual: no Claude agent harness session logs found" line, and box that line in red.
\u2022 IF YOU HAVE REAL USAGE EVIDENCE \u2014 a Teams message, an email, someone asking you a question the
  panel answered \u2014 use THAT instead and retitle the frame "someone asked, and the panel
  answered". Demand you can point at beats a defect report every time.
""")
    return s


# ------------------------------------------------------------------ challenges
def s_challenges(prs):
    s = blank(prs)
    y = header(s, "06  Reflection", "Challenges",
               "The six things that actually made this hard.")
    items = [
        ("The platform moved under me", MAGENTA,
         "Agent Host shipped mid-project. Collection got easier, interpretation got harder, "
         "and four assumptions stopped being true."),
        ("Three vendors, no shared vocabulary", AMBER,
         "Every harness spells the same concept differently \u2014 and all three ship weekly."),
        ("No ground truth to test against", PURPLE,
         "There is no reference answer for \u201chow many tokens did that session use\u201d."),
        ("One person, a moving target", BLUE,
         "Solo, against three harnesses that kept changing under the tests."),
        ("Failures that were completely silent", RED,
         "A port mismatch dropped 100% of telemetry with no error anywhere."),
        ("Performance in a place I didn\u2019t expect", TEAL,
         "A slow query froze the whole editor, not just the panel."),
    ]
    cw = Inches(5.83)
    for i, (title, col, body) in enumerate(items):
        x = MARGIN + (cw + Inches(0.23)) * (i % 2)
        yy = y + Inches(0.02) + Inches(1.5) * (i // 2)
        c = card(s, x, yy, cw, Inches(1.32), fill=WHITE, line=LINE)
        rect(s, x, yy, Inches(0.06), Inches(1.32), fill=col)
        tf = c.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.26)
        para(tf, title, size=13, color=NAVY, bold=True, first=True, space_after=5, line=1.14)
        para(tf, body, size=10.5, color=GRAY, space_after=0, line=1.22)
    footer(s, FOOT)
    notes(s, """
SCRIPT

Twenty-five seconds. Name them plainly. Do NOT apologise for any of them.

Say three, not six. Recommended: 1, 3, 4.

"Six things genuinely made this hard. I'll name three."

[#1 \u2014 already familiar, one sentence is enough]
"The platform moved under me. You've seen that one."

[#3 \u2014 the one engineers in the room will respect most]
"There was no ground truth to test against. There is no reference answer for 'how many tokens
did that session use' \u2014 nobody publishes one. So correctness had to be constructed. I built
fixtures per harness and wrote a test that asserts Copilot, Claude and Codex totals
simultaneously from one fixture. If any single convention drifts, that test fails."

[#4 \u2014 say it as judgement, not as shortfall]
"And I was one person against three harnesses that all shipped weekly. Which meant re-testing
assumptions constantly, and picking one surface to do properly rather than several half-way."

Then move. Don't editorialise.

\u2500\u2500\u2500 DETAIL FOR Q&A \u2500\u2500\u2500

\u2022 #5, silent failures: a port mismatch dropped 100% of telemetry with no error anywhere, and a
  thrown query error vanished into a console the user never opens \u2014 leaving "loading spans..."
  on screen forever. Both fixes were about making failure visible, and both mattered more than
  any feature.
\u2022 #6, performance: a slow query froze the whole editor, not just the panel. The two fixes that
  mattered were one SQL statement (ANALYZE, 8.4s \u2192 0.2s) and the physical column order in a
  table. Neither was where I'd have guessed.
\u2022 #2, no shared vocabulary: this is why every rule is written against observable shape rather
  than vendor name. Anything written as "if it's Codex, do X" was obsolete within a release.
""")
    return s


# ------------------------------------------------------------------ learning outcomes
def s_outcomes(prs):
    s = blank(prs)
    y = header(s, "06  Reflection", "Learning outcomes",
               "What I am taking out of this summer, in three buckets.")
    cols = [
        ("Engineering", "\u2699", BLUE,
         ["OpenTelemetry end to end \u2014 OTLP and the GenAI conventions",
          "Matching on shape, not on vendor name",
          "Normalization is the hard part \u2014 200 lines vs 2,600",
          "SQLite at speed \u2014 ANALYZE, indexes, column order",
          "Worker threads \u2014 never block the extension host",
          "Testing when there is no reference answer"]),
        ("Product & design", "\u25C6", PURPLE,
         ["Designing for two readers \u2014 a person and a model",
          "Writing for an LLM is its own discipline",
          "Choosing the unit of the product \u2014 session, not span",
          "Privacy as a design position, not a constraint",
          "Depth on one surface beat breadth across several"]),
        ("Working here", "\u2605", TEAL,
         ["Building on a platform that is still being built",
          "Reading an unfamiliar codebase to find what it emits",
          "Turning a written brief into a plan I could execute",
          "Knowing when to ship the narrower thing",
          "Asking upstream teams for what only they can fix"]),
    ]
    cw = Inches(3.83)
    for i, (title, icon, col, items) in enumerate(cols):
        x = MARGIN + (cw + Inches(0.2)) * i
        c = card(s, x, y, cw, Inches(4.52), fill=WHITE, line=LINE)
        rect(s, x, y, cw, Inches(0.055), fill=col)
        tf = c.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = tf.margin_right = Inches(0.26)
        tf.margin_top = Inches(0.26)
        para(tf, [(icon + "   ", dict(size=13, color=col, bold=True)),
                  (title, dict(size=15, color=NAVY, bold=True))],
             first=True, space_after=12)
        for it in items:
            para(tf, [("\u00b7  ", dict(color=col, bold=True)), (it, dict(color=GRAY))],
                 size=10.5, space_after=9, line=1.2)

    footer(s, FOOT)
    notes(s, """
SCRIPT

Twenty-five seconds. Do NOT read this slide \u2014 it's here so the audience can scan the breadth
while you say one sentence per column.

[open by naming the link to the previous slide, so the overlap reads as deliberate]
"Most of these came straight out of that last slide."

"On engineering \u2014 OpenTelemetry end to end, and the big one: the hard part was never collecting
the data, it was agreeing what it means. The receiver is two hundred lines. The engine that
interprets it is thirteen times that."

"On product \u2014 designing for two consumers at once, a person and a model, which turn out to want
very different things from the same data."

"And working here \u2014 building on a platform that is itself still being built, and learning that
when the plumbing gets solved for you, the value moves up a layer."

Then stop.

\u2500\u2500\u2500 DETAIL FOR Q&A \u2500\u2500\u2500

\u2022 "Matching on shape, not on vendor name": every filter written as "if it's Codex, do X" would
  have rotted within a release. The echo-detection rule is defined by shape instead \u2014 has
  content, no prompt, no round trip, no model-call span \u2014 which is why it survived three upstream
  releases untouched.
\u2022 Product: writing for an LLM means paging and budgeting instead of truncating, and saying
  "content capture was off" rather than returning an empty list, so the model can't infer what
  wasn't there.
\u2022 "The unit of the product": choosing the session rather than the span as the thing everything
  hangs off. That one decision is why the panel and the tools can share an engine.
\u2022 If asked which you'd want to go deeper on: the middle column. Designing for a model as a
  first-class reader is the part that felt genuinely new.
""")
    return s


# ------------------------------------------------------------------ thank you (photos)
def s_thanks_people(prs):
    s = blank(prs)
    bg(s, NAVY)
    accent_bar(s)
    tf = textbox(s, MARGIN, Inches(0.5), CONTENT_W, Inches(0.3))
    para(tf, "ACKNOWLEDGEMENTS", size=10.5, color=BLUE_L, bold=True, font=FONT_SB, first=True,
         space_after=0)
    tf = textbox(s, MARGIN, Inches(0.84), CONTENT_W, Inches(0.55))
    para(tf, "Special thank you", size=30, color=WHITE, font=FONT_LT, first=True, space_after=0)
    rect(s, MARGIN, Inches(1.54), Inches(1.05), Inches(0.035), fill=TEAL)

    people = [("@zhichli", "Framed the project"),
              ("@karthiknadig", "Early direction"),
              ("@kieferrm", "Early direction"),
              ("add name", "Manager / mentor"),
              ("add name", "Mentor / reviewer")]
    n = len(people)
    span = Inches(10.6)
    x0 = Inches(1.35)
    step = span / (n - 1)
    for i, (name, role) in enumerate(people):
        cx = x0 + step * i
        photo_slot(s, cx, Inches(2.86), Inches(1.42), "add\nphoto")
        tf = textbox(s, Emu(int(cx - Inches(1.7) / 2)), Inches(3.72), Inches(1.7), Inches(0.3))
        para(tf, name, size=11.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER, first=True,
             space_after=2)
        tf = textbox(s, Emu(int(cx - Inches(1.7) / 2)), Inches(4.02), Inches(1.7), Inches(0.5))
        para(tf, role.replace("\n", " "), size=9.5, color=RGBColor(0x93, 0xA6, 0xC4),
             align=PP_ALIGN.CENTER, first=True, space_after=0, line=1.15)

    teams = [("The VS Code agent-host team", "For microsoft/vscode#328529 \u2014 the session anchor "
              "span \u2014 and the chat.agentHost.otel.* surface this is built on.", PURPLE),
             ("The OpenTelemetry GenAI semconv WG", "gen_ai.* is the only reason a shared "
              "vocabulary exists to normalize toward. Without it this is three parsers and a prayer.",
              MAGENTA),
             ("The Copilot, Claude Code and Codex teams", "For emitting telemetry at all \u2014 three "
              "teams each independently decided this was worth instrumenting.", AMBER)]
    cw = Inches(3.83)
    for i, (name, body, col) in enumerate(teams):
        x = MARGIN + (cw + Inches(0.2)) * i
        c = card(s, x, Inches(4.78), cw, Inches(1.5), fill=RGBColor(0x1E, 0x33, 0x54),
                 line=RGBColor(0x30, 0x45, 0x6B))
        rect(s, x, Inches(4.78), cw, Inches(0.055), fill=col)
        tf = c.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.24)
        para(tf, name, size=12.5, color=col, bold=True, first=True, space_after=6, line=1.15)
        para(tf, body, size=10.5, color=RGBColor(0xB9, 0xC8, 0xE0), space_after=0, line=1.22)

    tf = textbox(s, MARGIN, Inches(6.46), Inches(11.89), Inches(0.4))
    para(tf, "\u2014 and everyone who dogfooded it and reported that their session list had "
             "261 rows in it.",
         size=12.5, color=RGBColor(0x8E, 0xA4, 0xC6), italic=True, first=True, space_after=0,
         align=PP_ALIGN.CENTER)

    notes(s, """
SCRIPT

Short, specific, sincere. Slow down here \u2014 it's the last real slide.

"Before I finish \u2014 a few thank-yous, and they're not perfunctory."

[name Zhichao first and properly]
"Zhichao framed this whole problem space before I got here. A lot of what I built is genuinely
an argument with ideas he'd already had, and that's the best possible starting point."

[then the mentors and manager \u2014 use their names]
"[names] \u2014 for the early direction, and for letting me change the shape of this when the
platform changed underneath it."

[then the three teams]
"The agent-host team, whose anchor span is the thing that made cross-harness session identity
tractable at all \u2014 even though, as you saw, it also handed me four new problems. Both are true."

"The OpenTelemetry GenAI working group. Their conventions are the only reason there's anything
to normalize toward. Without that, this is three parsers and a prayer."

"And the Copilot, Claude Code and Codex teams \u2014 for emitting this data at all. Three teams each
independently decided agent runs were worth instrumenting, and none of them had to."

[the last line \u2014 it usually gets a laugh]
"And everyone who dogfooded it and told me their session list had two hundred and sixty-one rows
in it."

\u2500\u2500\u2500 TO DO BEFORE PRESENTING \u2500\u2500\u2500

\u2022 Drop headshots into the five circles: right-click \u2192 Format Shape \u2192 Picture fill, or delete the
  circle and insert a picture over it.
\u2022 Replace the two "add name" slots with your manager and mentor.
\u2022 If you only have three people, delete the extras and re-space the remaining circles evenly.
""")
    return s


# ------------------------------------------------------------------ close + contact
def s_close_contact(prs):
    s = blank(prs)
    bg(s, NAVY)
    rect(s, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill=BLUE)
    xs = [10.05, 10.62, 11.19, 11.76, 12.33]
    hs = [1.2, 2.1, 0.8, 1.7, 1.05]
    cols = [BLUE, TEAL, PURPLE, BLUE_L, MAGENTA]
    for x, h, c in zip(xs, hs, cols):
        rect(s, Inches(x), Inches(4.02 - h), Inches(0.30), Inches(h), fill=c,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.35)

    tf = textbox(s, Inches(1.05), Inches(1.12), Inches(8.6), Inches(0.3))
    para(tf, "THANK YOU", size=12, color=BLUE_L, bold=True, font=FONT_SB, first=True,
         space_after=0)
    tf = textbox(s, Inches(1.0), Inches(1.56), Inches(8.8), Inches(2.4))
    para(tf, [("The data layer is solved.", dict(color=WHITE))],
         size=38, font=FONT_LT, first=True, space_after=6, line=1.1)
    para(tf, [("Agent Insights is about ", dict(color=RGBColor(0x93, 0xA6, 0xC4))),
              ("everything after.", dict(color=BLUE_L, bold=True))],
         size=38, font=FONT_LT, space_after=0, line=1.1)
    rect(s, Inches(1.05), Inches(3.72), Inches(1.5), Inches(0.045), fill=TEAL)

    tf = textbox(s, Inches(1.05), Inches(4.02), Inches(8.4), Inches(0.6))
    para(tf, "Three agents, three vocabularies, one machine \u2014 and now one definition of a "
             "session, a turn, a model call, a tool call and a token.",
         size=14, color=RGBColor(0xB9, 0xC8, 0xE0), first=True, space_after=0, line=1.28)

    cards_ = [("\u2709", "michiisai22@gmail.com", BLUE_L),
              ("in", "linkedin.com/in/michiisai", TEAL),
              ("\u2691", "github.com/michiisai/agent-insights", PURPLE)]
    x = Inches(1.05)
    for icon, val, col in cards_:
        c = card(s, x, Inches(5.06), Inches(3.55), Inches(0.72),
                 fill=RGBColor(0x1E, 0x33, 0x54), line=RGBColor(0x30, 0x45, 0x6B))
        tf = c.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.24)
        para(tf, [(icon + "   ", dict(size=12, color=col, bold=True)),
                  (val, dict(size=10.5, color=RGBColor(0xC3, 0xD1, 0xE6), font=MONO))],
             first=True, space_after=0)
        x += Inches(3.75)

    tf = textbox(s, Inches(1.05), Inches(6.1), Inches(9.0), Inches(0.4))
    para(tf, "Questions? \u2014 and there\u2019s an appendix with the numbers if you want detail.",
         size=13, color=RGBColor(0x7B, 0x92, 0xB8), italic=True, first=True, space_after=0)

    notes(s, """
SCRIPT

Close, then stop talking.

[deliver the two lines, and let the second one sit for a beat]
"The data layer is solved."
"Agent Insights is about everything after."

[pause \u2014 this is the conclusion the last twenty minutes exists to earn]

[point at the grey line \u2014 this is the bookend]
"Three agents, three vocabularies, one machine \u2014 and now one definition of a session, a turn, a
model call, a tool call and a token."

[say "and now" with a little weight; that's the whole talk in two words]

"Happy to take questions. My contact details are up here, and there's an appendix with all the
numbers if anyone wants to go deeper on the normalization."

Leave this slide up for the whole Q&A \u2014 that's what the contact cards are for.

\u2500\u2500\u2500 NOTES \u2500\u2500\u2500

\u2022 This is the FIRST time the audience hears the "data layer" line \u2014 it only works because they
  now know what the data layer is and what "after" means. Don't rush it.
\u2022 The title slide said "three agents, three vocabularies, one machine" and left it hanging.
  This closes it. Same seven opening words, deliberately.

\u2500\u2500\u2500 LIKELY QUESTIONS \u2500\u2500\u2500

\u2022 "Why LM tools and not MCP?" \u2014 first-party: no separate server process, they appear as
  #-references, and the bundled skill lets the agent route without being named. MCP was the
  obvious first guess; LM tools turned out strictly better inside VS Code.
\u2022 "Why sql.js and not better-sqlite3?" \u2014 no native module to rebuild per platform and per
  Electron ABI. The cost is a WASM boundary, which is exactly why it lives on a worker thread.
\u2022 "How do you know your numbers are right?" \u2014 522 assertions across 15 suites, including one
  that asserts Copilot, Claude and Codex token totals simultaneously from a single fixture.
\u2022 "Privacy?" \u2014 loopback only, content capture opt-in, threat model in the repo. There's an
  appendix slide with the full data-flow diagram.
\u2022 "What would you do with another month?" \u2014 cross-process OTel inside VS Code itself, and
  per-model pricing so the status bar shows spend rather than usage.
\u2022 "What was the hardest bug?" \u2014 Codex response.completed records carry no trace id, so joining
  by trace matched nothing and no Codex transcript reported a single token. Fixed by joining on
  conversation id instead.
""")
    return s


# ------------------------------------------------------------------ appendix: data flow
def s_appendix_dataflow(prs):
    s = blank(prs)
    y = header(s, "Appendix \u00b7 privacy", "Where the data goes",
               "Every arrow is on your machine except the last one \u2014 and that one only "
               "moves when you ask it to.")

    bw = Inches(1.98)
    gap = Inches(0.5)
    bh = Inches(0.9)
    cy = y + Inches(0.82)
    xs = [MARGIN + (bw + gap) * i for i in range(5)]

    # trust boundary around the three local stages
    bx0 = xs[1] - Inches(0.22)
    bx1 = xs[3] + bw + Inches(0.22)
    band = rect(s, bx0, y + Inches(0.28), bx1 - bx0, Inches(2.72),
                fill=WASH, line=RGBColor(0xA8, 0xB8, 0xD4))
    band.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    band.shadow.inherit = False
    tf = textbox(s, bx0 + Inches(0.16), y + Inches(0.36), Inches(5.0), Inches(0.28))
    para(tf, "YOUR MACHINE  \u00b7  the Agent Insights extension", size=9, color=BLUE, bold=True,
         font=FONT_SB, first=True, space_after=0)

    nodes = [("AI agents", "Copilot \u00b7 Claude Code \u00b7 Codex",
              RGBColor(0xFD, 0xF6, 0xE3), AMBER),
             ("Local OTLP receiver", "loopback HTTP only",
              RGBColor(0xEC, 0xF6, 0xEC), GREEN),
             ("SQLite on disk", "extension storage",
              RGBColor(0xEA, 0xF3, 0xFB), BLUE),
             ("Panel + 12 LM tools", "renderer and agent surface",
              RGBColor(0xEC, 0xF6, 0xEC), GREEN),
             ("Model provider", "via Copilot Chat",
              RGBColor(0xFD, 0xF6, 0xE3), AMBER)]
    for x, (label, sub, fill, line_c) in zip(xs, nodes):
        sh = card(s, x, cy, bw, bh, fill=fill, line=line_c)
        tf = sh.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.12)
        para(tf, label, size=11, color=NAVY, bold=True, align=PP_ALIGN.CENTER,
             first=True, space_after=3, line=1.1)
        para(tf, sub, size=8.5, color=GRAY, align=PP_ALIGN.CENTER, space_after=0, line=1.12)

    labels = ["1. telemetry", "2. store", "3. read", "5. on request only"]
    for i, lab in enumerate(labels):
        ax = xs[i] + bw + Inches(0.05)
        aw = gap - Inches(0.1)
        last = (i == len(labels) - 1)
        rect(s, ax, cy + Inches(0.38), aw, Inches(0.15),
             fill=MAGENTA if last else RGBColor(0xC8, 0xC6, 0xC4),
             shape=MSO_SHAPE.RIGHT_ARROW)
        tf = textbox(s, ax - Inches(0.42), cy + bh + Inches(0.04), aw + Inches(0.84), Inches(0.3))
        para(tf, lab, size=8.5, color=MAGENTA if last else GRAY, bold=last,
             align=PP_ALIGN.CENTER, first=True, space_after=0)

    # the user, inside the boundary
    ux = xs[2] - Inches(0.3)
    u = card(s, ux, cy + Inches(1.42), Inches(2.6), Inches(0.62),
             fill=RGBColor(0xFD, 0xF6, 0xE3), line=AMBER)
    fill_shape_text(u, [("You", dict(size=11.5, color=NAVY, bold=True,
                                     align=PP_ALIGN.CENTER, space_after=0))])
    tf = textbox(s, ux + Inches(2.72), cy + Inches(1.5), Inches(3.4), Inches(0.5))
    para(tf, "4. browse \u00b7 clear \u00b7 opt in to\ncontent capture", size=9, color=GRAY,
         first=True, space_after=0, line=1.18)

    band2 = card(s, MARGIN, Inches(5.62), CONTENT_W, Inches(1.0), fill=NAVY, line=NAVY)
    rect(s, MARGIN, Inches(5.62), Inches(0.075), Inches(1.0), fill=MAGENTA)
    tf = band2.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.36)
    para(tf, [("The only arrow that leaves is the one you trigger. ",
               dict(color=WHITE, bold=True)),
              ("Asking Copilot Chat a question sends that answer to the model provider \u2014 the "
               "same trust boundary as any other chat message. No collector, no account, no "
               "background upload, and prompt and response content is opt-in.",
               dict(color=RGBColor(0xB9, 0xC8, 0xE0)))],
         size=12, first=True, space_after=0, line=1.24)

    footer(s, FOOT + "  \u00b7  appendix")
    notes(s, """
APPENDIX \u2014 not in the running order. Jump here for any privacy or "where does the data go"
question, which is the most likely serious question you will get.

Trace it left to right with your hand:
  1. The agents export OTLP to a port on localhost. Loopback only \u2014 not reachable from another
     machine.
  2. The receiver writes to a SQLite file inside the extension's own storage directory.
  3. The panel and the twelve LM tools read from that file. Nothing else does.
  4. You control it \u2014 browse it, clear it, and content capture (prompts and replies) is off
     unless you turn it on.
  5. The one magenta arrow: Copilot Chat answering you. Same trust boundary as any other chat
     message you send.

The strong version of the claim: "there is no collector to stand up, no account, and no
background upload. If you unplug the network, everything on the previous slides still works."

Full threat model is in the repo at docs/threat-model.md.
""")
    return s
