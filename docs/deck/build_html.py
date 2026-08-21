"""Generate a self-contained HTML preview deck from the built .pptx + rendered PNGs."""
import base64
import html
import re
import sys
from pathlib import Path

from pptx import Presentation

PPTX = Path(sys.argv[1])
PNG_DIR = Path(sys.argv[2])
OUT = Path(sys.argv[3])


def read_notes(prs):
    out = []
    for slide in prs.slides:
        text = ""
        if slide.has_notes_slide:
            text = slide.notes_slide.notes_text_frame.text or ""
        out.append(text.strip())
    return out


def main():
    prs = Presentation(str(PPTX))
    notes = read_notes(prs)
    images = []
    for i in range(1, len(notes) + 1):
        for name in ("Slide%d.PNG" % i, "Slide%d.png" % i):
            p = PNG_DIR / name
            if p.exists():
                images.append(base64.b64encode(p.read_bytes()).decode("ascii"))
                break
        else:
            raise SystemExit("missing render for slide %d" % i)

    cards = []
    for i, (img, note) in enumerate(zip(images, notes), start=1):
        timing = ""
        m = re.match(r"\[([^\]]+)\]", note)
        if m:
            timing = m.group(1)
            note = note[m.end():].strip()
        body = html.escape(note).replace("\n\n", "</p><p>").replace("\n", "<br>")
        cards.append(
            '<section class="slide" id="s{n}" data-n="{n}">'
            '<div class="stage"><img alt="Slide {n}" src="data:image/png;base64,{img}"></div>'
            '<aside class="notes"><header><span class="num">{n} / {tot}</span>'
            '{timing}</header><div class="body"><p>{body}</p></div></aside>'
            '</section>'.format(
                n=i, tot=len(images), img=img,
                timing=('<span class="timing">%s</span>' % html.escape(timing)) if timing else "",
                body=body or "<em>No speaker notes.</em>")
        )

    tpl = TEMPLATE.replace("{{SLIDES}}", "\n".join(cards)).replace("{{TOTAL}}", str(len(images)))
    OUT.write_text(tpl, encoding="utf-8")
    print("wrote %s  (%.1f MB, %d slides)" % (OUT, OUT.stat().st_size / 1e6, len(images)))


TEMPLATE = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>Agent Insights &mdash; Final Presentation</title>
<style>
  :root{
    --navy:#142 43e; --navy:#14243e; --navy2:#1e3354; --blue:#0078d4; --blueL:#50e6ff;
    --ink:#1b1a19; --gray:#605e5c; --line:#e1dfdd; --wash:#f7f7f9; --white:#fff;
  }
  *{box-sizing:border-box}
  html,body{margin:0;padding:0}
  body{
    font-family:"Segoe UI","Segoe UI Variable",system-ui,-apple-system,sans-serif;
    background:#eef1f6;color:var(--ink);-webkit-font-smoothing:antialiased;
  }
  header.top{
    position:sticky;top:0;z-index:50;background:var(--navy);color:#fff;
    display:flex;align-items:center;gap:18px;padding:10px 22px;
    box-shadow:0 1px 0 rgba(255,255,255,.06),0 6px 18px rgba(0,0,0,.18);
  }
  header.top .bar{position:absolute;left:0;right:0;top:0;height:3px;
    background:linear-gradient(90deg,#0078d4 0 25%,#00b7c3 25% 50%,#8661c5 50% 75%,#e3008c 75% 100%)}
  header.top h1{font-size:15px;font-weight:600;margin:0;letter-spacing:.2px}
  header.top .sub{font-size:12px;color:#93a6c4;margin-left:-8px}
  header.top .spacer{flex:1}
  header.top button, header.top a.btn{
    background:rgba(255,255,255,.08);color:#dce6f5;border:1px solid rgba(255,255,255,.16);
    border-radius:14px;padding:5px 13px;font:inherit;font-size:12px;cursor:pointer;
    text-decoration:none;transition:background .12s;
  }
  header.top button:hover,header.top a.btn:hover{background:rgba(255,255,255,.18)}
  header.top button.on{background:var(--blue);border-color:var(--blue);color:#fff}
  header.top .pos{font-size:12px;color:#93a6c4;font-variant-numeric:tabular-nums;min-width:62px;text-align:right}

  main{max-width:1500px;margin:0 auto;padding:26px 22px 90px}
  .slide{
    display:grid;grid-template-columns:minmax(0,1fr) 380px;gap:22px;
    margin:0 0 34px;scroll-margin-top:70px;
  }
  body.nonotes .slide{grid-template-columns:minmax(0,1fr)}
  body.nonotes .notes{display:none}
  .stage{background:#fff;border-radius:10px;overflow:hidden;align-self:start;
    box-shadow:0 1px 2px rgba(0,0,0,.10),0 10px 28px rgba(16,24,40,.14)}
  .stage img{display:block;width:100%;height:auto}
  .notes{background:#fff;border:1px solid var(--line);border-radius:10px;
    display:flex;flex-direction:column;overflow:hidden;align-self:start;
    max-height:calc(100vh - 120px);position:sticky;top:70px}
  .notes header{display:flex;align-items:center;gap:10px;padding:11px 16px;
    border-bottom:1px solid var(--line);background:var(--wash)}
  .notes .num{font-size:11px;font-weight:700;color:var(--blue);letter-spacing:.4px}
  .notes .timing{font-size:11px;font-weight:600;color:#fff;background:var(--navy2);
    border-radius:10px;padding:2px 9px;font-variant-numeric:tabular-nums}
  .notes .body{padding:14px 16px 18px;overflow:auto;font-size:13.5px;line-height:1.55;color:#3b3a39}
  .notes .body p{margin:0 0 11px;white-space:normal}
  .notes .body p:last-child{margin-bottom:0}
  .notes .body em{color:var(--gray)}

  /* presentation mode */
  body.present{overflow:hidden;background:#0b1220}
  body.present header.top{opacity:0;pointer-events:none}
  body.present main{max-width:none;margin:0;padding:0}
  body.present .slide{display:none;margin:0}
  body.present .slide.active{display:block;height:100vh}
  body.present .stage{border-radius:0;box-shadow:none;height:100vh;align-self:stretch;
    display:flex;align-items:center;justify-content:center;background:#0b1220}
  body.present .stage img{max-height:100vh;width:auto;max-width:100%}
  body.present .notes{display:none}
  body.present.shownotes .slide.active{display:grid;grid-template-columns:minmax(0,1fr) 400px;height:100vh}
  body.present.shownotes .notes{display:flex;border-radius:0;border:0;border-left:1px solid #23324a;
    background:#0f1829;color:#c9d6ea;max-height:100vh;height:100vh;position:static}
  body.present.shownotes .notes header{background:#16233a;border-color:#23324a}
  body.present.shownotes .notes .body{color:#c3d1e6;font-size:15px}
  body.present.shownotes .stage{height:100vh;min-width:0}
  body.present.shownotes .stage img{max-width:100%;max-height:100vh;width:auto}
  .hint{position:fixed;bottom:14px;left:50%;transform:translateX(-50%);
    background:rgba(20,36,62,.92);color:#c3d1e6;font-size:12px;padding:7px 15px;
    border-radius:16px;z-index:60;pointer-events:none;opacity:0;transition:opacity .3s}
  body.present .hint{opacity:1}
  @media print{
    header.top,.hint{display:none}
    main{max-width:none;padding:0}
    .slide{display:block;page-break-after:always;margin:0}
    .notes{border:0;margin-top:8px;position:static;max-height:none}
    .stage{box-shadow:none;border-radius:0}
  }
</style>
</head>
<body>
<header class="top">
  <div class="bar"></div>
  <h1>Agent Insights</h1>
  <span class="sub">Internship Final Presentation &middot; {{TOTAL}} slides</span>
  <span class="spacer"></span>
  <span class="pos" id="pos">1 / {{TOTAL}}</span>
  <button id="toggleNotes" class="on">Notes</button>
  <button id="present">Present &nbsp;&#9654;</button>
  <a class="btn" href="Agent-Insights-Final-Presentation.pptx" download>Download .pptx</a>
</header>
<main id="main">
{{SLIDES}}
</main>
<div class="hint">&larr; &rarr; / Space navigate &nbsp;&middot;&nbsp; N toggles notes &nbsp;&middot;&nbsp; Esc exits</div>
<script>
(function(){
  var slides=[].slice.call(document.querySelectorAll('.slide'));
  var total=slides.length, idx=0;
  var pos=document.getElementById('pos');
  var body=document.body;

  function setIdx(n){
    idx=Math.max(0,Math.min(total-1,n));
    pos.textContent=(idx+1)+' / '+total;
    if(body.classList.contains('present')){
      slides.forEach(function(s,i){s.classList.toggle('active',i===idx);});
    }
  }
  function enterPresent(){
    body.classList.add('present');
    if(document.documentElement.requestFullscreen){
      document.documentElement.requestFullscreen().catch(function(){});
    }
    setIdx(idx);
  }
  function exitPresent(){
    body.classList.remove('present');
    if(document.fullscreenElement && document.exitFullscreen){document.exitFullscreen();}
    slides[idx].scrollIntoView({block:'start'});
  }
  document.getElementById('present').addEventListener('click',enterPresent);
  document.getElementById('toggleNotes').addEventListener('click',function(){
    if(body.classList.contains('present')){body.classList.toggle('shownotes');}
    else{body.classList.toggle('nonotes');this.classList.toggle('on');}
  });

  document.addEventListener('keydown',function(e){
    var p=body.classList.contains('present');
    if(e.key==='Escape'&&p){exitPresent();return;}
    if(e.key==='n'||e.key==='N'){
      if(p){body.classList.toggle('shownotes');}
      else{body.classList.toggle('nonotes');document.getElementById('toggleNotes').classList.toggle('on');}
      return;
    }
    if(e.key==='f'||e.key==='F'){if(!p)enterPresent();return;}
    if(!p) return;
    if(e.key==='ArrowRight'||e.key==='ArrowDown'||e.key===' '||e.key==='PageDown'){e.preventDefault();setIdx(idx+1);}
    if(e.key==='ArrowLeft'||e.key==='ArrowUp'||e.key==='PageUp'){e.preventDefault();setIdx(idx-1);}
    if(e.key==='Home'){setIdx(0);} if(e.key==='End'){setIdx(total-1);}
  });
  document.addEventListener('click',function(e){
    if(!body.classList.contains('present'))return;
    if(e.target.closest('header.top'))return;
    setIdx(idx+(e.clientX<window.innerWidth*0.25?-1:1));
  });

  // track position while scrolling in review mode
  var io=new IntersectionObserver(function(entries){
    if(body.classList.contains('present'))return;
    entries.forEach(function(en){
      if(en.isIntersecting){
        idx=slides.indexOf(en.target);
        pos.textContent=(idx+1)+' / '+total;
      }
    });
  },{threshold:.5});
  slides.forEach(function(s){io.observe(s);});
  setIdx(0);
})();
</script>
</body>
</html>
"""

if __name__ == "__main__":
    main()
