"""Slides 1-16: title, problem, approach."""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from theme import *



# ------------------------------------------------------------------ 1 title
def s_title(prs):
    s = blank(prs)
    bg(s, NAVY)
    rect(s, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill=BLUE)
    # decorative signal bars, top right
    xs = [10.05, 10.62, 11.19, 11.76, 12.33]
    hs = [1.55, 2.35, 0.95, 1.95, 1.25]
    cols = [BLUE, TEAL, PURPLE, BLUE_L, MAGENTA]
    for x, h, c in zip(xs, hs, cols):
        rect(s, Inches(x), Inches(4.62 - h), Inches(0.30), Inches(h), fill=c,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.35)

    tf = textbox(s, Inches(1.05), Inches(1.62), Inches(9.2), Inches(0.3))
    para(tf, "MICROSOFT  ·  VISUAL STUDIO CODE  ·  SEPT 2026", size=11.5,
         color=BLUE_L, bold=True, font=FONT_SB, first=True, space_after=0)

    tf = textbox(s, Inches(1.0), Inches(2.12), Inches(9.4), Inches(1.3))
    para(tf, "Agent Insights", size=68, color=WHITE, font=FONT_LT, first=True, space_after=0)

    rect(s, Inches(1.05), Inches(3.42), Inches(1.5), Inches(0.045), fill=TEAL)

    tf = textbox(s, Inches(1.05), Inches(3.78), Inches(8.6), Inches(1.1))
    para(tf, "OpenTelemetry for AI coding agents, inside VS Code", size=23,
         color=RGBColor(0xC9, 0xD8, 0xEE), font=FONT_LT, first=True, space_after=8, line=1.15)
    para(tf, "Copilot, Claude Code and Codex \u2014 three agents, three vocabularies, one machine.",
         size=13.5, color=RGBColor(0x93, 0xA6, 0xC4), space_after=0, line=1.25)

    tf = textbox(s, Inches(1.05), Inches(5.72), Inches(10.5), Inches(1.0))
    para(tf, [("Michelle Ma", dict(bold=True, color=WHITE)),
              ("   |   Software Engineer Intern", {})],
         size=13.5, color=RGBColor(0x93, 0xA6, 0xC4), first=True, space_after=6)
    para(tf, "github.com/michiisai/agent-insights",
         size=12, color=BLUE_L, font=MONO, space_after=0)

    notes(s, """
SCRIPT

[Cold open. Do NOT advance yet. Have your status bar visible, or describe it.]

"Before I say anything about the project \u2014 this is my VS Code status bar."

   $(broadcast)  \u219312.4K   42% cached   \u21913.1K

"That's what my coding agents cost me today. Twelve thousand tokens in, three thousand out,
forty-two percent of the input served from cache."

"Nobody built me a dashboard for that. That number came out of telemetry the agents were already
emitting into the void."

[beat]

"I'm Michelle, and I spent my internship building Agent Insights: a VS Code extension that picks
that telemetry up and turns it into something you \u2014 and your agent \u2014 can actually read."

"Two short recorded demos partway through, and questions at the end."

[advance]

\u2500\u2500\u2500 NOTES \u2500\u2500\u2500

\u2022 This is the ONLY place in the talk you define the project. Make it land, then stop \u2014 the next
  fifteen minutes are the explanation, and the problem section works far better if the room
  feels the pain before it hears the solution.
\u2022 If the status bar isn't showing real numbers on the day, just say the figures out loud. The
  point is that a number exists and nothing explains it.
\u2022 There's an appendix with all the numbers if anyone wants to go deeper in Q&A.
""")
    return s


# ------------------------------------------------------------------ 2 thesis
def s_thesis(prs):
    s = blank(prs)
    bg(s, WASH_2)
    accent_bar(s)
    rect(s, Inches(0), Inches(0.075), Inches(0.16), SLIDE_H, fill=NAVY)

    tf = textbox(s, MARGIN, Inches(0.95), Inches(11.6), Inches(2.0))
    para(tf, [("The data layer for agent observability is ", dict(color=NAVY)),
              ("solved", dict(color=GREEN, bold=True)), (".", dict(color=NAVY))],
         size=38, color=NAVY, font=FONT_LT, first=True, space_after=6, line=1.12)
    para(tf, [("Everything ", dict(color=NAVY)),
              ("after you turn it on", dict(color=MAGENTA, bold=True)),
              (" is not.", dict(color=NAVY))],
         size=38, color=NAVY, font=FONT_LT, space_after=0, line=1.12)

    tf = textbox(s, MARGIN, Inches(2.62), Inches(11.4), Inches(0.6))
    para(tf, "Copilot, Claude Code and Codex all emit OpenTelemetry today. "
             "Nothing on your machine reads it.",
         size=15.5, color=GRAY, first=True, space_after=0)

    # the one-liner
    sh = card(s, MARGIN, Inches(3.5), Inches(11.89), Inches(1.62), fill=WHITE, line=LINE)
    rect(s, MARGIN, Inches(3.5), Inches(0.075), Inches(1.62), fill=BLUE)
    tf = sh.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.42)
    tf.margin_right = Inches(0.36)
    para(tf, "THE ONE-LINER", size=10, color=BLUE, bold=True, font=FONT_SB, first=True,
         space_after=8)
    para(tf, [("Agent Insights collects the OpenTelemetry your coding agents "
               "already emit, and turns it into ", {}),
              ("readable agent sessions", dict(bold=True, color=NAVY)),
              (" \u2014 for you in a panel, and for your agent through "
               "language-model tools.", {})],
         size=17.5, color=INK, space_after=0, line=1.28)

    # three quick stat chips
    stats = [("3", "agent harnesses\nnormalized", BLUE),
             ("12", "language-model tools\n+ 1 chat skill", PURPLE),
             ("100%", "local \u2014 nothing leaves the machine", TEAL)]
    x = MARGIN
    for val, label, col in stats:
        c = card(s, x, Inches(5.42), Inches(3.83), Inches(1.28), fill=WHITE, line=LINE)
        rect(s, x, Inches(5.42), Inches(3.83), Inches(0.055), fill=col)
        tf = c.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.28)
        para(tf, val, size=28, color=col, font=FONT_LT, bold=True, first=True, space_after=1)
        para(tf, label.replace("\n", "  "), size=11.5, color=GRAY, space_after=0, line=1.15)
        x += Inches(4.03)

    footer(s, FOOT)
    notes(s, """
This is the thesis of the whole talk, and I'll come back to it at the end.

OpenTelemetry won. Every serious coding agent \u2014 Copilot, Claude Code, Codex \u2014 can emit it
today, out of the box, with a settings toggle. That fight is over.

What nobody solved is the part after the toggle. You turn it on, and it goes to... where?
A collector you have to stand up? A Jaeger instance that has never heard of the word "turn"?

So: Agent Insights. Read the one-liner off the slide \u2014 it's the sentence I want people to
repeat afterwards. Emphasise "readable agent sessions" and "and for your agent", because
those are the two halves of the product.
""")
    return s


# ------------------------------------------------------------------ 3 agenda
def s_agenda(prs):
    s = blank(prs)
    y = header(s, "Agenda", "Twenty minutes, six beats",
               "Slides through the middle; the demo is live on pre-captured data.")
    items = [
        ("01", "Problem", "Why agent runs are the least observable thing you run", "4 min", BLUE),
        ("02", "The twist", "VS Code shipped Agent Host mid-project", "1 min", MAGENTA),
        ("03", "Approach", "Architecture, and normalizing three vocabularies", "6 min", PURPLE),
        ("04", "Demo", "Two scenarios, both ending on the loop closing", "4\u00bd min", TEAL),
        ("05", "Impact", "What shipped, and where it goes next", "2 min", GREEN),
        ("06", "Reflection", "What made it hard, and what I'd tell you", "2 min", AMBER),
    ]
    top = Inches(1.98)
    h = Inches(0.68)
    gap = Inches(0.115)
    for i, (num, title, sub, mins, col) in enumerate(items):
        yy = top + (h + gap) * i
        c = card(s, MARGIN, yy, Inches(11.89), h, fill=WHITE, line=LINE)
        rect(s, MARGIN, yy, Inches(0.07), h, fill=col)
        tf = textbox(s, MARGIN + Inches(0.34), yy + Inches(0.15), Inches(0.7), Inches(0.4))
        para(tf, num, size=19, color=col, font=FONT_LT, bold=True, first=True, space_after=0)
        tf = textbox(s, MARGIN + Inches(1.05), yy + Inches(0.12), Inches(2.4), Inches(0.4))
        para(tf, title, size=16, color=NAVY, bold=True, first=True, space_after=0)
        tf = textbox(s, MARGIN + Inches(3.5), yy + Inches(0.18), Inches(6.6), Inches(0.4))
        para(tf, sub, size=12.5, color=GRAY, first=True, space_after=0)
        ch = card(s, Inches(11.55), yy + Inches(0.17), Inches(0.92), Inches(0.34),
                  fill=WASH, line=LINE, radius=0.5)
        fill_shape_text(ch, [(mins, dict(size=10.5, color=GRAY, bold=True,
                                         align=PP_ALIGN.CENTER, space_after=0))])
    footer(s, FOOT)
    notes(s, """
Move fast here \u2014 20 seconds, don't read every line.

Point at 02: "There's a plot twist in the middle. The platform I was building on shipped a
feature mid-project that changed which problem I was solving. That's the most interesting
part of the talk and I'm not going to hide it."

Point at 04: "Demo runs on data I captured earlier this week \u2014 I'm not generating telemetry
live on stage."
""")
    return s


# ------------------------------------------------------------------ dividers
def divider(prs, num, title, sub, mins=None, color=BLUE):
    s = blank(prs)
    bg(s, NAVY)
    rect(s, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill=color)
    tf = textbox(s, Inches(1.15), Inches(2.32), Inches(2.0), Inches(1.4))
    para(tf, num, size=88, color=color, font=FONT_LT, bold=True, first=True, space_after=0)
    rect(s, Inches(3.05), Inches(2.42), Inches(0.03), Inches(2.0), fill=RGBColor(0x3A, 0x4E, 0x71))
    tf = textbox(s, Inches(3.55), Inches(2.5), Inches(8.6), Inches(1.9))
    para(tf, title, size=42, color=WHITE, font=FONT_LT, first=True, space_after=12, line=1.1)
    para(tf, sub, size=15.5, color=RGBColor(0x93, 0xA6, 0xC4), space_after=0, line=1.3)
    return s


# ------------------------------------------------------------------ 5 user problem
def s_problem_user(prs):
    s = blank(prs)
    y = header(s, "01  Problem \u00b7 the user", "Five questions you cannot answer about your own agent runs")
    qs = [
        ("\u201cWhat did I even ask it to do?\u201d", "Six sessions, two harnesses, one morning.", BLUE),
        ("\u201cWhy did that take four minutes?\u201d", "Model? Tool call? Permission prompt? Retry?", PURPLE),
        ("\u201cWhy did it go wrong?\u201d", "Agents rarely fail loudly. They misread you quietly.", MAGENTA),
        ("\u201cWhat is this costing me?\u201d", "Tokens, cache reads, cache writes, subagents.", AMBER),
        ("\u201cWas Claude better than Copilot here?\u201d", "A gut feeling, and zero evidence.", TEAL),
    ]
    col_w = Inches(3.83)
    positions = [(MARGIN, y + Inches(0.05)), (MARGIN + Inches(4.03), y + Inches(0.05)),
                 (MARGIN + Inches(8.06), y + Inches(0.05)),
                 (MARGIN, y + Inches(1.62)), (MARGIN + Inches(4.03), y + Inches(1.62))]
    for (title, sub, col), (x, yy) in zip(qs, positions):
        c = card(s, x, yy, col_w, Inches(1.42), fill=WHITE, line=LINE)
        rect(s, x, yy, col_w, Inches(0.055), fill=col)
        tf = c.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.24)
        para(tf, title, size=14, color=NAVY, bold=True, first=True, space_after=6, line=1.15)
        para(tf, sub, size=11.5, color=GRAY, space_after=0, line=1.24)

    # punchline banner in the empty 6th cell
    x = MARGIN + Inches(8.06)
    c = card(s, x, y + Inches(1.62), col_w, Inches(1.42), fill=NAVY, line=NAVY)
    tf = c.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.24)
    para(tf, "And the loop that should exist doesn\u2019t.", size=13.5, color=BLUE_L,
         bold=True, first=True, space_after=6, line=1.15)
    para(tf, "The agent in your editor reads structured data all day \u2014 "
             "and has no access to its own history.",
         size=11.5, color=RGBColor(0xB9, 0xC8, 0xE0), space_after=0, line=1.24)

    band = card(s, MARGIN, Inches(5.66), Inches(11.89), Inches(1.0), fill=BLUE_XL, line=BLUE_XL)
    tf = band.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.36)
    tf.margin_right = Inches(0.36)
    para(tf, [("An agent session is the most expensive thing a developer runs all day, "
               "and the ", {}),
              ("least observable thing on their machine.", dict(bold=True, color=BLUE))],
         size=17, color=NAVY, first=True, space_after=0, line=1.2)

    footer(s, FOOT)
    notes(s, """
SCRIPT

"Here are five questions I could not answer about my own agent runs. Not exotic questions \u2014
these are the ones you ask every single day."

[read three only \u2014 1, 3 and 5]

"What did I even ask it to do? I'd have six sessions open across two harnesses by lunchtime, and
no idea which one was the one that actually worked."

"Why did it go wrong? \u2014 and this is the one that justifies the whole project. Agents almost
never throw an exception. They misread you quietly on turn two, and then spend eight more turns
confidently building on the wrong assumption. That failure is semantic. There's no error, no
latency spike, no non-zero exit code. Every metric you'd normally collect says everything is
fine."

"And was Claude actually better than Copilot for this task? I had a gut feeling and exactly zero
evidence."

[point at the navy card]
"There's also a loop that should exist and doesn't. The agent sitting in your editor reads
structured data all day \u2014 and it has no access to its own history."

[the band \u2014 say it slowly, then pause]
"An agent session is the most expensive thing a developer runs all day, and the least observable
thing on their machine."

\u2500\u2500\u2500 DETAIL FOR Q&A \u2500\u2500\u2500

\u2022 The comparison to make if it helps: you can profile a 200-millisecond function call in a dozen
  tools. You cannot profile the four-minute, forty-cent, fourteen-tool-call conversation that
  just rewrote your auth layer.
\u2022 "Expensive" is meant literally as well as figuratively \u2014 tokens cost money, and the runs are
  long enough that you context-switch away and lose the thread.
\u2022 The navy card is the setup for the language-model tools later. Flag it here, don't explain it.
""")
    return s


# ------------------------------------------------------------------ 6 why logs fail
def s_problem_logs(prs):
    s = blank(prs)
    y = header(s, "01  Problem \u00b7 why not just read the logs",
               "Four reasons the data that already exists doesn\u2019t answer any of that",
               "Every item below is something the codebase had to work around.")
    items = [
        ("1", "The in-box view has gaps under Agent Host", BLUE,
         "VS Code ships Agent Debug Logs \u2014 agent-host sessions don\u2019t reliably appear in it."),
        ("2", "Append-only text, not a queryable model", PURPLE,
         "A transcript says what was said. Duration and causality are what a log throws away."),
        ("3", "OTel exists \u2014 and nothing consumes it", MAGENTA,
         "Turn it on and you need an endpoint. A collector plus Jaeger, or nothing."),
        ("4", "Generic viewers show spans; humans think in conversations", AMBER,
         "Jaeger has never heard of a \u201cturn\u201d. It will show you 261 traces and name none of them."),
    ]
    top = y + Inches(0.06)
    h = Inches(1.06)
    gap = Inches(0.15)
    for i, (num, title, col, body) in enumerate(items):
        yy = top + (h + gap) * i
        c = card(s, MARGIN, yy, Inches(11.89), h, fill=WHITE, line=LINE)
        rect(s, MARGIN, yy, Inches(0.07), h, fill=col)
        nb = card(s, MARGIN + Inches(0.3), yy + Inches(0.29), Inches(0.48), Inches(0.48),
                  fill=col, line=col, radius=0.5)
        fill_shape_text(nb, [(num, dict(size=15, color=WHITE, bold=True,
                                        align=PP_ALIGN.CENTER, space_after=0))])
        tf = textbox(s, MARGIN + Inches(1.0), yy + Inches(0.16), Inches(10.6), Inches(0.85))
        para(tf, title, size=14.5, color=NAVY, bold=True, first=True, space_after=4)
        para(tf, body, size=11.5, color=GRAY, space_after=0, line=1.22)
    footer(s, FOOT)
    notes(s, """
SCRIPT

"So why not just read the logs? Four reasons, and I hit all four."

[#1]
"First \u2014 and I want to get ahead of this, because half of you are already thinking it \u2014 VS Code
does ship Agent Debug Logs. But there's an open issue right now that Claude harness sessions
don't show up in them at all. The in-box view hasn't caught up with the agent host yet. And
underneath that there's no shared format anyway: Claude writes JSONL transcripts, Codex has its
own session files, Copilot keeps history in the workbench. Three formats, three locations, zero
joins."

[#2 \u2014 the strongest technical point, slow down]
"Second, and this is the real one: a transcript is append-only text. It tells you what was said.
It cannot tell you that the third tool call took forty-one seconds, or that turn five spent
ninety percent of its input on cache reads, or that two of those twelve tool calls were retries.
Duration, causality, parent-child structure \u2014 that's exactly what a chronological log destroys
by construction."

[#3 \u2014 say the privacy line deliberately]
"Third, OpenTelemetry does exist, and nothing consumes it. Turn it on and it asks you for an
endpoint. Your options are stand up a collector and a Jaeger instance \u2014 heavyweight, off-machine,
and privacy-hostile for data that is literally your prompts and your source code \u2014 or nothing.
Most people pick nothing."

[#4 \u2014 the segue]
"And fourth, even if you do that, generic viewers show you spans. Jaeger has never heard of a
turn. It will happily show you two hundred and sixty-one traces, and not one of them is labelled
'the conversation where I asked it to fix the flaky test.'"

Then: "hold onto that 261 \u2014 it comes back."

\u2500\u2500\u2500 DETAIL FOR Q&A \u2500\u2500\u2500

\u2022 Be generous about the debug logs. It's a gap in something shipping fast, not a failing, and it
  will get fixed. Don't name the issue number on stage.
\u2022 The durable version of #1: even once the debug logs work, they're per-harness text, not a
  queryable cross-agent model. That gap doesn't close.
\u2022 Privacy is why this became local-first, and it comes back on the Impact slide \u2014 flag it here
  as a design position rather than a limitation.
\u2022 261 is a real number from one day of my own use. It comes back in the Agent Host section.
""")
    return s


# ------------------------------------------------------------------ otel primer
def s_otel(prs):
    s = blank(prs)
    y = header(s, "01  Problem \u00b7 the standard", "OpenTelemetry, in thirty seconds",
               "An open standard for recording what a program did while it ran \u2014 "
               "and all three agents already speak it.")

    # ---- left: a trace is a tree of spans
    lw = Inches(7.0)
    panel = card(s, MARGIN, y + Inches(0.04), lw, Inches(4.42), fill=WHITE, line=LINE)
    rect(s, MARGIN, y + Inches(0.04), lw, Inches(0.055), fill=BLUE)
    tf = textbox(s, MARGIN + Inches(0.3), y + Inches(0.26), lw - Inches(0.6), Inches(0.3))
    para(tf, "A TRACE IS A TREE OF SPANS", size=10.5, color=BLUE, bold=True, font=FONT_SB,
         first=True, space_after=0)

    rows = [(0, "chat turn", 0.00, 2.40, BLUE),
            (1, "model call", 0.04, 1.90, PURPLE),
            (2, "http POST /messages", 0.10, 1.70, TEAL),
            (1, "tool call \u00b7 read_file", 1.96, 0.26, AMBER),
            (1, "tool call \u00b7 edit_file", 2.22, 0.18, AMBER)]
    bar_x0 = MARGIN + Inches(3.42)
    bar_max = Inches(2.6)
    total = 2.40
    ry = y + Inches(0.72)
    for depth, name, start, dur, col in rows:
        x = MARGIN + Inches(0.34) + Inches(0.26) * depth
        tf = textbox(s, x, ry, Inches(3.0), Inches(0.28))
        para(tf, name, size=10.5, color=NAVY if depth == 0 else GRAY, bold=(depth == 0),
             font=MONO, first=True, space_after=0)
        off = Emu(int(bar_max * (start / total)))
        w = max(Emu(int(bar_max * (dur / total))), Inches(0.14))
        rect(s, bar_x0 + off, ry + Inches(0.045), w, Inches(0.19), fill=col,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
        tf = textbox(s, bar_x0 + bar_max + Inches(0.12), ry, Inches(0.9), Inches(0.28))
        para(tf, "%.2f s" % dur, size=9.5, color=GRAY_L, font=MONO, first=True, space_after=0)
        ry += Inches(0.42)

    note = card(s, MARGIN + Inches(0.34), ry + Inches(0.16), lw - Inches(0.68), Inches(1.28),
                fill=WASH_2, line=WASH_2)
    tf = note.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.26)
    para(tf, [("One span ", dict(color=BLUE, bold=True)),
              ("= one unit of work, with a name, a start time, a duration, a parent, "
               "and a bag of attributes.", dict(color=NAVY))],
         size=11.5, first=True, space_after=6, line=1.22)
    para(tf, [("One trace ", dict(color=BLUE, bold=True)),
              ("= every span that shares a trace id. That is the whole data model.",
               dict(color=NAVY))],
         size=11.5, space_after=0, line=1.22)

    # ---- right: signals + semconv
    rx = MARGIN + lw + Inches(0.23)
    rw = Inches(4.66)
    signals = [("Traces", BLUE, "What happened, in what order, and how long each part took"),
               ("Metrics", TEAL, "Counters and histograms sampled over time"),
               ("Logs", PURPLE, "Timestamped events, attachable to a span")]
    sy = y + Inches(0.04)
    for name, col, body in signals:
        c = card(s, rx, sy, rw, Inches(0.92), fill=WHITE, line=LINE)
        rect(s, rx, sy, Inches(0.06), Inches(0.92), fill=col)
        tf = textbox(s, rx + Inches(0.3), sy + Inches(0.13), rw - Inches(0.5), Inches(0.3))
        para(tf, name, size=12.5, color=NAVY, bold=True, first=True, space_after=2)
        para(tf, body, size=10.5, color=GRAY, space_after=0, line=1.18)
        sy += Inches(1.02)

    cs = card(s, rx, sy + Inches(0.12), rw, Inches(1.32), fill=NAVY, line=NAVY)
    rect(s, rx, sy + Inches(0.12), rw, Inches(0.055), fill=MAGENTA)
    tf = cs.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.28)
    para(tf, "SEMANTIC CONVENTIONS", size=9.5, color=MAGENTA, bold=True, font=FONT_SB,
         first=True, space_after=6)
    para(tf, [("gen_ai.*", dict(font=MONO, color=WHITE, bold=True)),
              ("  \u2014 agreed names so \u201cwhich model\u201d and \u201chow many tokens\u201d mean the same thing "
               "everywhere.", dict(color=RGBColor(0xB9, 0xC8, 0xE0)))],
         size=11, space_after=0, line=1.22)

    footer(s, FOOT)
    notes(s, """
SCRIPT

Thirty-five seconds. Pure setup \u2014 you are buying comprehension for the next fifteen minutes.
Don't go faster than the room.

"Quick thirty seconds on OpenTelemetry, because everything after this uses two words from it."

"It's an open standard for recording what a program did while it ran. Most modern observability
is built on it. And the useful thing here is that all three of these agents already emit it \u2014
the data exists, nobody was reading it."

[trace the diagram with your hand, top to bottom]
"A span is one unit of work. It has a name, a start time, a duration, a parent, and a bag of
attributes. A trace is every span that shares an id. That is genuinely the entire data model."

[right side, one breath]
"Three kinds of signal \u2014 traces, metrics, logs."

[the magenta box \u2014 this is the bridge to the whole talk, land it deliberately]
"And there are semantic conventions. Agreed attribute names, so that 'which model' or 'how many
tokens' means the same thing no matter who emitted it."

[beat]

"That last part is where this entire project lives \u2014 because in practice, the three harnesses
do not agree."

Then pivot hard: "so the data exists, and it's in a standard format. Here's what you still
can't answer."

\u2500\u2500\u2500 NOTES \u2500\u2500\u2500

\u2022 If the room is all observability people, say the one-liner and move through fast. But if there
  is a single PM or manager present, KEEP THIS SLIDE \u2014 everything after it uses "span" and
  "trace" without apologising.
\u2022 The waterfall is deliberately the same visual grammar as the Codex span cascade two slides
  later, so the audience already knows how to read that one.
\u2022 If asked what OTLP is: the wire protocol \u2014 how the data gets from the agent to anything that
  wants to receive it.
""")
    return s


# ------------------------------------------------------------------ 7 span vs turn
def s_span_vs_turn(prs):
    s = blank(prs)
    y = header(s, "01  Problem \u00b7 the crux",
               "The gap between a span and a turn is the entire product",
               "Two ways the raw signal refuses to line up with how a person thinks.")

    # LEFT: codex nesting
    tf = textbox(s, MARGIN, y + Inches(0.02), Inches(6.6), Inches(0.3))
    para(tf, "ONE CODEX MODEL CALL, AS OTEL SEES IT", size=10.5, color=MAGENTA, bold=True,
         font=FONT_SB, first=True, space_after=0)

    spans = ["run_sampling_request", "try_run_sampling_request", "stream_request",
             "model_client.stream_responses_api", "responses.stream_request"]
    top = y + Inches(0.42)
    for i, name in enumerate(spans):
        x = MARGIN + Inches(0.34 * i)
        w = Inches(5.9) - Inches(0.34 * i)
        col = MAGENTA if i == 0 else RGBColor(0xE9, 0xE7, 0xE6)
        fg = WHITE if i == 0 else GRAY
        sh = card(s, x, top + Inches(0.46 * i), w, Inches(0.36), fill=col, line=col, radius=0.12)
        fill_shape_text(sh, [(name, dict(size=11, color=fg, bold=(i == 0), font=MONO,
                                         space_after=0))], pad=Inches(0.16))
    tf = textbox(s, MARGIN, top + Inches(2.42), Inches(6.3), Inches(0.5))
    para(tf, [("5 spans. ", dict(bold=True, color=NAVY)),
              ("1 model call. All five report the same token count. "
               "A generic viewer calls that five operations.", dict(color=GRAY))],
         size=12, first=True, space_after=0, line=1.25)

    # RIGHT: attribute chaos
    rx = Inches(7.62)
    tf = textbox(s, rx, y + Inches(0.02), Inches(5.0), Inches(0.3))
    para(tf, "AND NOBODY AGREES ON THE NAMES", size=10.5, color=AMBER, bold=True,
         font=FONT_SB, first=True, space_after=0)

    code_block(s, rx, y + Inches(0.42), Inches(4.99), Inches(2.26), [
        [("// \u201cinput tokens\u201d, three ways", dict(color=RGBColor(0x7A, 0x93, 0xB8)))],
        [("gen_ai.usage.input_tokens", dict(color=BLUE_L)),
         ("      OTel semconv", dict(color=RGBColor(0x7A, 0x93, 0xB8)))],
        [("llm.usage.prompt_tokens", dict(color=BLUE_L)),
         ("        Copilot legacy", dict(color=RGBColor(0x7A, 0x93, 0xB8)))],
        [("input_tokens", dict(color=BLUE_L)),
         ("                   Claude Code", dict(color=RGBColor(0x7A, 0x93, 0xB8)))],
        "",
        [("// \u201cmodel\u201d, four ways", dict(color=RGBColor(0x7A, 0x93, 0xB8)))],
        [("gen_ai.request.model  gen_ai.response.model", dict(color=BLUE_L))],
        [("llm.model             model", dict(color=BLUE_L))],
    ], size=11)

    tf = textbox(s, rx, y + Inches(2.82), Inches(4.9), Inches(0.5))
    para(tf, "Ask a generic viewer \u201chow many tokens today?\u201d and it answers with whichever "
             "third of your data happened to use the key it knows.",
         size=12, color=GRAY, first=True, space_after=0, line=1.25)

    band = card(s, MARGIN, Inches(5.72), Inches(11.89), Inches(0.92), fill=NAVY, line=NAVY)
    tf = band.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.36)
    para(tf, [("Worse: some of the conversation isn\u2019t in the traces at all. ",
               dict(bold=True, color=BLUE_L)),
              ("Claude puts your prompt on a span and the model\u2019s reply in a log record. "
               "Codex reports content only as logs \u2014 and strips the payload out of "
               "every assistant message, so it never exports the model\u2019s words at all.",
               dict(color=RGBColor(0xC3, 0xD1, 0xE6)))],
         size=13, first=True, space_after=0, line=1.24)

    footer(s, FOOT)
    notes(s, """
SCRIPT

This is the most important slide in the Problem section. Slow down.

"So here's the crux, in two pictures."

[left \u2014 walk down the nesting with your hand]
"This is one model call, as OpenTelemetry sees it on Codex. Five spans. Not five operations \u2014
one. Every single one of them reports the same token count."

"So if you naively count model calls, you're five times wrong on Codex and exactly right on
Copilot. Which means any cross-agent comparison is garbage before you've started."

[right]
"And nobody agrees on the names. Three different spellings of 'input tokens'. Four of 'model'.
That's not a nitpick \u2014 it's a hard blocker on every aggregate question you'd want to ask."

"Ask a generic viewer 'how many tokens today' and it answers with whichever third of your data
happened to use the key it knows about."

[the band \u2014 this is the one that surprised me most]
"And it gets worse. Two of the three split the two halves of a conversation across two different
signals \u2014 your prompt on a span, the model's reply in a log record. So a tool that reads traces
sees half a conversation, and a tool that reads logs sees a different half."

"And Codex strips the model's words out entirely before export. It never sends them at all."

[land it]
"That gap \u2014 between a span, which is what you're given, and a turn, which is what you actually
want \u2014 that gap is the entire product."

\u2500\u2500\u2500 DETAIL FOR Q&A \u2500\u2500\u2500

\u2022 On Codex only the outermost of the five spans is counted as the model call. The other four
  are internal layers of the same request.
\u2022 The Codex assistant-text gap is the one thing nothing downstream can fix \u2014 it's an upstream
  ask, and it's #1 on the roadmap slide.
\u2022 If asked why the vendors differ: none of them are wrong. The GenAI conventions describe how to
  WRITE this data; nothing describes how to read it back as a conversation. Each team made a
  reasonable local choice.
""")
    return s


# ------------------------------------------------------------------ 8 agent host
def s_agent_host(prs):
    s = blank(prs)
    y = header(s, "02  The twist \u00b7 mid-project",
               "VS Code shipped Agent Host \u2014 and moved the hard part",
               "Claude Code and Codex became native VS Code sessions. Collection got easier. "
               "Interpretation got harder.")

    # gave
    gx, gw = MARGIN, Inches(4.36)
    c = card(s, gx, y + Inches(0.04), gw, Inches(2.9), fill=WHITE, line=LINE)
    rect(s, gx, y + Inches(0.04), gw, Inches(0.055), fill=GREEN)
    tf = c.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = Inches(0.26)
    tf.margin_top = Inches(0.22)
    para(tf, "WHAT IT GAVE ME", size=10.5, color=GREEN, bold=True, font=FONT_SB, first=True,
         space_after=9)
    for t in ["One setting to turn telemetry on",
              "An anchor span linking the host to the agent\u2019s own trace",
              "A conversation id stamped onto that trace",
              "A session title, for free"]:
        para(tf, [("\u2713  ", dict(color=GREEN, bold=True)), (t, dict(color=INK))],
             size=11.5, space_after=7, line=1.22)

    # broke
    bx = MARGIN + Inches(4.62)
    bw = Inches(7.27)
    c = card(s, bx, y + Inches(0.04), bw, Inches(2.9), fill=WHITE, line=LINE)
    rect(s, bx, y + Inches(0.04), bw, Inches(0.055), fill=MAGENTA)
    tf = c.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = Inches(0.26)
    tf.margin_top = Inches(0.22)
    para(tf, "AND THE FOUR PROBLEMS IT CREATED", size=10.5, color=MAGENTA, bold=True,
         font=FONT_SB, first=True, space_after=9)
    rows = [
        ("It relabels every session.",
         "MAX(service_name) returns vscode-agent-host \u2014 it just sorts last."),
        ("The in-box debug view hasn\u2019t caught up.",
         "Agent Debug Logs don\u2019t reliably show harness sessions."),
        ("The anchor only anchors the first trace.",
         "Later turns carry no conversation id at all."),
        ("A conversation id isn\u2019t evidence of a conversation.",
         "The host mints one when a chat is created, not when it\u2019s used."),
    ]
    for i, (t, sub) in enumerate(rows):
        para(tf, [("%d  " % (i + 1), dict(color=MAGENTA, bold=True)),
                  (t, dict(color=NAVY, bold=True)), ("  " + sub, dict(color=GRAY))],
             size=11, space_after=6, line=1.2)

    # stat band
    band = card(s, MARGIN, Inches(5.12), Inches(11.89), Inches(1.52), fill=NAVY, line=NAVY)
    rect(s, MARGIN, Inches(5.12), Inches(0.075), Inches(1.52), fill=MAGENTA)
    tf = textbox(s, MARGIN + Inches(0.42), Inches(5.4), Inches(3.9), Inches(1.1))
    para(tf, [("261", dict(size=44, color=MAGENTA, bold=True, font=FONT_LT)),
              ("  vs  ", dict(size=20, color=RGBColor(0x7A, 0x93, 0xB8))),
              ("6", dict(size=44, color=GREEN, bold=True, font=FONT_LT))],
         first=True, space_after=2)
    para(tf, "phantom sessions vs real ones, in one day of use",
         size=11.5, color=RGBColor(0x93, 0xA6, 0xC4), space_after=0)
    tf = textbox(s, Inches(5.34), Inches(5.4), Inches(7.16), Inches(1.0), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "Open a Codex thread and never type in it, and you still get a \u201csession\u201d. "
             "The real work disappears into the noise.",
         size=12.5, color=RGBColor(0xC3, 0xD1, 0xE6), first=True, space_after=6, line=1.22)
    para(tf, [("The takeaway: the platform moving under me didn\u2019t invalidate the project \u2014 "
               "it moved the hard part from ", dict(color=RGBColor(0x93, 0xA6, 0xC4))),
              ("\u201ccollect the data\u201d", dict(color=WHITE, bold=True)),
              (" to ", dict(color=RGBColor(0x93, 0xA6, 0xC4))),
              ("\u201cdecide what the data means\u201d.", dict(color=WHITE, bold=True))],
         size=12.5, space_after=0, line=1.22)

    footer(s, FOOT)
    notes(s, """
SCRIPT

"About halfway through my internship, VS Code shipped Agent Host. Claude Code and Codex stopped
being separate terminal tools and became native sessions inside the editor. This is what happens
when you build on a platform that is itself being built."

[left column \u2014 15 seconds, don't labour it]
"And genuinely, it helped. One setting to turn telemetry on. An anchor span that links the host
to the agent's own trace. A conversation id stamped onto it. Even a session title, for free. On
paper, that solves session identity for nothing."

[right column]
"In practice it created four problems. My favourite bug of the summer is the first one: every
single native session was suddenly labelled 'the agent host'. Not for any deep reason \u2014 I was
taking a MAX over the service name, and 'vscode-agent-host' happens to sort after 'claude-code'
alphabetically."

"The second one is the one that made this project matter more rather than less: the in-box debug
view hasn't caught up with the host yet."

"And the last two are about that conversation id. It only anchors the first trace \u2014 later turns
carry nothing. And the host mints one when a chat is created, not when it's used."

[the number \u2014 let it land]
"Which gets you this. Two hundred and sixty-one phantom sessions burying six real ones, in a
single day of my own use. Open a Codex thread, never type in it, and you still get a session."

[the takeaway \u2014 say it deliberately]
"The takeaway is the one I'd keep: the platform moving under me didn't invalidate the project.
It moved the hard part from collecting the data to deciding what the data means."

Then: "and it isn't just me hitting that gap."

\u2500\u2500\u2500 DETAIL FOR Q&A \u2500\u2500\u2500

\u2022 Fix for the relabelling: exclude host spans explicitly from span counts and service labels,
  rather than trusting an aggregate over service name.
\u2022 Fix for the phantom sessions: keep a session only if it shows real agent activity, a captured
  user prompt, or a host-assigned title. A trace id fallback alone mints one session per
  housekeeping trace.
\u2022 An empty Codex thread produces roughly 37 spans of pure startup \u2014 that's what "phantom" means
  here.
\u2022 This is the setup for the "platform moved under you" learning at the end. Don't spend it here.
""")
    return s


# ------------------------------------------------------------------ 10 the table
def s_diff_table(prs):
    s = blank(prs)
    y = header(s, "03  Approach \u00b7 the input", "Three harnesses, three vocabularies",
               "Every row is something that had to be handled or the number on screen would be wrong.")
    rows = [
        ["", "Copilot", "Claude Code", "Codex"],
        ["Service name", "github-copilot", "claude-code", "codex-app-server"],
        ["One model call =", "one span", "one span", "five nested spans"],
        ["One tool call =", "execute_tool%", "claude_code.tool", "handle_tool_call"],
        ["Token attributes", "gen_ai.usage.*", "bare input_tokens", "inside a log record"],
        ["Cache tokens", "subset of input", "added on top of input", "subset of input"],
        ["User prompt", "on the span", "on the span", "in a log record"],
        ["Assistant text", "on the span", "in a log record", "never exported"],
        ["Session identity", "chat_session_id", "conversation.id", "often nothing"],
        ["Log timestamps", "usable as sent", "usable as sent", "always zero"],
    ]
    tbl = add_table(s, rows, MARGIN, y + Inches(0.02), Inches(11.89),
                    col_w=[1.55, 2.85, 3.1, 3.2], row_h=Inches(0.335), header_h=Inches(0.38))
    style_table(tbl, size=10, header_size=11)
    # colour the vendor headers
    for c, col in [(1, BLUE), (2, AMBER), (3, MAGENTA)]:
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = col

    band = card(s, MARGIN, Inches(5.5), Inches(11.89), Inches(1.3), fill=BLUE_XL, line=BLUE_XL)
    tf = band.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.34)
    tf.margin_top = tf.margin_bottom = Inches(0.12)
    para(tf, [("Counting a tool call is not obvious. ", dict(bold=True, color=BLUE)),
              ("On Claude, matching the whole subtree trebles the count; matching too narrowly "
               "drops every permission-denied call.", dict(color=NAVY))],
         size=12.5, first=True, space_after=6, line=1.22)
    para(tf, [("Counting a model call is not obvious. ", dict(bold=True, color=BLUE)),
              ("On Codex, five spans report the same call. ", dict(color=NAVY)),
              ("Reading a conversation is not obvious. ", dict(bold=True, color=BLUE)),
              ("Two of three split it across two signals; one never sends the model\u2019s half.",
               dict(color=NAVY))],
         size=12.5, space_after=0, line=1.22)

    footer(s, FOOT)
    notes(s, """
SCRIPT

This is the most persuasive slide in the deck and you get it almost for free \u2014 do NOT read it
row by row.

"I'm not going to read this to you. I want you to look at the shape of it. Four columns, ten
rows, and almost nothing lines up."

[three seconds of silence \u2014 let them scan]

"Three rows to actually notice."

"Cache tokens. Copilot and Codex follow the OpenTelemetry convention, where cache reads are a
subset of your input tokens \u2014 already counted. Anthropic's are additive, on top. Same idea,
opposite arithmetic. Get that wrong and you under-report Claude's prompt tokens by three or four
times."

"Assistant text \u2014 what the model actually said. Three different answers, and one of them is
'never'. Codex strips it before export."

"And one model call. On Codex that's five nested spans. Count them naively and every Codex
session looks five times busier than it is."

[the band]
"So: counting a tool call is not obvious. Counting a model call is not obvious. Reading a
conversation is not obvious. That's why this is an engineering problem and not a rendering
problem."

\u2500\u2500\u2500 DETAIL FOR Q&A \u2500\u2500\u2500

\u2022 Claude tool calls: matching the whole subtree trebles the count, because each tool call has
  child spans. Matching only the .execution child silently drops every permission-denied call \u2014
  which are exactly the ones you most want to see. The fix is an exact match on claude_code.tool.
\u2022 Log timestamps: Codex sends a literal zero and puts the real clock in observedTimeUnixNano.
  A plain COALESCE only falls through on NULL, so the zero won \u2014 every Codex log landed at 1970,
  sorted last, and vanished from any time-windowed query.
\u2022 Session identity: Copilot uses its own chat_session_id, Claude uses the GenAI conversation id
  or a session id, and Codex frequently carries neither \u2014 it has to be recovered from logs.
\u2022 If someone asks why not just normalise at ingest: because the vendors ship weekly. Anything
  written as "if it's Codex, do X" is a dated cheque. The rules are written against observable
  shape instead.
""")
    return s
