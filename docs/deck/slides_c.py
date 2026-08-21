"""Slides: demo, impact, next steps, learnings, acknowledgements, appendix."""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from theme import *



def _storyboard(s, y, steps, col):
    x = MARGIN
    w = Inches(2.26)
    gap = Inches(0.145)
    for i, (title, body) in enumerate(steps):
        c = card(s, x, y, w, Inches(2.5), fill=WHITE, line=LINE)
        rect(s, x, y, w, Inches(0.055), fill=col)
        nb = card(s, x + Inches(0.2), y + Inches(0.24), Inches(0.4), Inches(0.4),
                  fill=col, line=col, radius=0.5)
        fill_shape_text(nb, [(str(i + 1), dict(size=13, color=WHITE, bold=True,
                                               align=PP_ALIGN.CENTER, space_after=0))])
        tf = textbox(s, x + Inches(0.2), y + Inches(0.76), w - Inches(0.4), Inches(1.6))
        para(tf, title, size=12, color=NAVY, bold=True, first=True, space_after=5, line=1.15)
        para(tf, body, size=10.5, color=GRAY, space_after=0, line=1.22)
        x += w + gap
        if i < len(steps) - 1:
            tri = rect(s, x - gap - Inches(0.02), y + Inches(0.34), Inches(0.14), Inches(0.2),
                       fill=RGBColor(0xC8, 0xC6, 0xC4), shape=MSO_SHAPE.RIGHT_ARROW)


def _video_slide(s, y, color, watch, payoff_big, payoff_small, clip_hint):
    """Left: a 16:9 slot for the screen recording. Right: what to watch for, then the payoff."""
    vw = Inches(7.9)
    vh = Inches(4.44)
    slot = rect(s, MARGIN, y + Inches(0.04), vw, vh, fill=WASH_2, line=RGBColor(0xA8, 0xB8, 0xD4))
    slot.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    slot.shadow.inherit = False
    tf = slot.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, "\u25B6", size=30, color=RGBColor(0x9F, 0xB2, 0xCE), align=PP_ALIGN.CENTER,
         first=True, space_after=8)
    para(tf, "drop the screen recording here", size=13, color=GRAY,
         align=PP_ALIGN.CENTER, space_after=4)
    para(tf, clip_hint, size=10.5, color=GRAY_L, align=PP_ALIGN.CENTER, space_after=0)

    rx = MARGIN + vw + Inches(0.22)
    rw = Inches(3.77)
    tf = textbox(s, rx, y + Inches(0.04), rw, Inches(0.3))
    para(tf, "WHAT TO WATCH FOR", size=10, color=color, bold=True, font=FONT_SB,
         first=True, space_after=0)

    ry = y + Inches(0.42)
    for i, (title, body) in enumerate(watch):
        c = card(s, rx, ry, rw, Inches(1.02), fill=WHITE, line=LINE)
        rect(s, rx, ry, Inches(0.055), Inches(1.02), fill=color)
        nb = card(s, rx + Inches(0.22), ry + Inches(0.15), Inches(0.32), Inches(0.32),
                  fill=color, line=color, radius=0.5)
        fill_shape_text(nb, [(str(i + 1), dict(size=10.5, color=WHITE, bold=True,
                                               align=PP_ALIGN.CENTER, space_after=0))])
        tf = textbox(s, rx + Inches(0.62), ry + Inches(0.13), rw - Inches(0.8), Inches(0.36))
        para(tf, title, size=11.5, color=NAVY, bold=True, first=True, space_after=0)
        tf = textbox(s, rx + Inches(0.22), ry + Inches(0.52), rw - Inches(0.42), Inches(0.44))
        para(tf, body, size=9.5, color=GRAY, first=True, space_after=0, line=1.16)
        ry += Inches(1.1)

    band = card(s, rx, ry + Inches(0.04), rw, Inches(1.02), fill=NAVY, line=NAVY)
    rect(s, rx, ry + Inches(0.04), rw, Inches(0.05), fill=color)
    tf = band.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.24)
    para(tf, payoff_big, size=12.5, color=WHITE, bold=True, first=True, space_after=4, line=1.15)
    para(tf, payoff_small, size=9.5, color=RGBColor(0x93, 0xA6, 0xC4), space_after=0, line=1.18)


def s_demo1(prs):
    s = blank(prs)
    y = header(s, "04  Demo \u00b7 scenario one", "\u201cWhy did that session go wrong?\u201d",
               "The semantic-failure demo \u2014 the one no existing tool can do at all.")
    _video_slide(
        s, y, BLUE,
        [("The session list",
          "Every row assembled out of spans that never contained the word \u201csession\u201d. "
          "Three different agent badges."),
         ("Turn 5, in words",
          "Duration and tokens spike. Open the transcript and the failure is a misread "
          "requirement, not an exception."),
         ("Ask the agent about itself",
          "\u201c+ chat\u201d pre-stages #agentSession. It reads its own telemetry and "
          "answers with deeplinks.")],
        "The agent just diagnosed the agent.",
        "And its answer is clickable \u2014 the whole loop closing, on this machine.",
        "~90 seconds  \u00b7  1920\u00d71080  \u00b7  no audio")
    footer(s, FOOT)
    notes(s, """
SCRIPT \u2014 recorded clip, you narrate live over it.

Do NOT read the right-hand column. It's there so the audience knows where to look while you talk.

[before you start the clip]
"First scenario: why did that session go wrong. This is a recording \u2014 I wanted to show you the
real thing without betting the talk on a live agent run."

[press play, then talk over it]

"So this is the session list. Every row here was assembled out of spans that never once
contained the word 'session'. Three different agents, one list, and they're directly
comparable \u2014 same columns, same meanings."

"I'm picking a run that took too long. And what you get is a turn-by-turn timeline \u2014 you can see
duration and tokens both spike right here, on turn five."

"Open that turn, and this is the part no other tool gets you to. The failure isn't an exception.
There's no error, no crash. It misread a requirement on turn five and then spent eight more turns
confidently building on it. That's a semantic failure, and it's only visible in the words."

"Now watch this \u2014 I click a sentence in the transcript, and it lands on the exact span that
produced it. The conversation and the trace are the same object, viewed two ways."

"And then the loop closes. This button pre-stages a reference to the session in chat. I send it,
and the agent reads its own telemetry through the tools and tells me what happened \u2014 with links
back into the panel, so every claim it makes is checkable."

[as the clip ends, land the line]
"The agent just diagnosed the agent. And its answer is clickable."

Do not improvise past that line.

\u2500\u2500\u2500 RECORDING SETUP \u2500\u2500\u2500

\u2022 Target ~90 seconds, 1920\u00d71080, no audio \u2014 you are the narration.
\u2022 Bump VS Code's font size before recording; text that's fine on your monitor is unreadable from
  the back of a room.
\u2022 Speed up any wait longer than ~2 seconds by 2-4\u00d7 rather than cutting it \u2014 it still reads as
  real, but nobody watches a spinner.
\u2022 In PowerPoint: set the video to Start Automatically and turn Loop until Stopped OFF.
\u2022 Have the panel already open when the clip starts. Don't spend seconds on navigation.
""")
    return s


def s_demo2(prs):
    s = blank(prs)
    y = header(s, "04  Demo \u00b7 scenario two", "\u201cWhich agent should I have used?\u201d",
               "The normalization payoff \u2014 it looks boring and it is the hardest thing "
               "in the project.")
    _video_slide(
        s, y, PURPLE,
        [("Status bar \u2192 Home",
          "\u219312.4K  42% cached  \u21913.1K. Tokens by model, cache read vs write, tool error "
          "rates, background calls."),
         ("Three harnesses, one task",
          "Three session rows side by side. Same columns, same meaning \u2014 that is the "
          "entire contribution."),
         ("It states its own limits",
          "Told that call counts are not like-for-like across harnesses, it "
          "describes rather than ranks.")],
        "\u201cGetting these three rows comparable is 90% of this project.\u201d",
        "One definition of a session, a turn, a call and a token \u2014 across three vocabularies.",
        "~75 seconds  \u00b7  1920\u00d71080  \u00b7  no audio")
    footer(s, FOOT)
    notes(s, """
SCRIPT \u2014 recorded clip, you narrate live over it.

This one looks boring on purpose. Your job is to say why it isn't.

[before you start]
"Second scenario, and it's the opposite of the first \u2014 this one looks completely unremarkable,
and it's the hardest thing in the project."

[press play, then talk over it]

"Starting in the status bar. That's today's tokens and cache hit rate, updating as runs land.
One click opens the panel."

"This is Home. Tokens by model, cache reads versus cache writes, tool call error rates. And
those two cache columns mean different arithmetic depending on which agent produced them \u2014 this
chart has already reconciled that. You just can't see it happening."

[the three rows \u2014 SLOW DOWN, this is the payoff of the whole technical section]
"And here's the moment. Same task, three harnesses, three rows side by side. Same columns. Same
meanings."

"Getting these three rows comparable is ninety percent of this project."

"Codex reports one model call as five nested spans. Claude puts the model's reply in a log
record. All three spell 'input tokens' differently. And none of that is visible here, because
that's the point."

"Last thing \u2014 when I ask it to compare them, watch what it does. It tells me the call counts
aren't like-for-like across harnesses instead of just ranking them. A tool that knows the limits
of its own numbers is more credible than one that doesn't."

\u2500\u2500\u2500 IF YOU'RE RUNNING LONG \u2500\u2500\u2500

This is the clip to shorten. Trim it to the three rows and the caveat; drop the Home tour
entirely. The comparison is the only part that has to survive.

\u2500\u2500\u2500 RECORDING SETUP \u2500\u2500\u2500

\u2022 Target ~75 seconds, 1920\u00d71080, no audio.
\u2022 Have all three sessions already captured and visible \u2014 don't record yourself scrolling to
  find them.
\u2022 The comparison rows should be on screen for at least 8 seconds. That's the shot people will
  remember.
""")
    return s


def s_impact(prs):
    s = blank(prs)
    y = header(s, "05  Impact", "What is true now that wasn\u2019t before")
    items = [
        ("Agent sessions are a first-class, queryable object", BLUE,
         "On your own machine, across three harnesses, with one definition of a turn, a model "
         "call, a tool call and a token."),
        ("Semantic failures are diagnosable", MAGENTA,
         "\u201cIt misunderstood me on turn 5\u201d is now something you can point at."),
        ("The agent can reason about itself", GREEN,
         "\u201cWhy was that slow\u201d and \u201ccompare these two runs\u201d are answerable in the editor \u2014 "
         "with clickable answers."),
        ("Cost is visible continuously", AMBER,
         "A live status-bar baseline, built from telemetry that was otherwise thrown away."),
        ("Local-first and private by construction", TEAL,
         "Nothing leaves the machine \u2014 which matters precisely because the data is prompts "
         "and source code."),
        ("It survives real use", PURPLE,
         "Multi-window handover, bounded retention, and a worker thread so a heavy query "
         "can\u2019t freeze the editor."),
    ]
    cw = Inches(5.83)
    for i, (title, col, body) in enumerate(items):
        x = MARGIN + (cw + Inches(0.23)) * (i % 2)
        yy = y + Inches(0.02) + Inches(1.16) * (i // 2)
        c = card(s, x, yy, cw, Inches(1.0), fill=WHITE, line=LINE)
        rect(s, x, yy, Inches(0.06), Inches(1.0), fill=col)
        tf = c.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.28)
        tf.margin_right = Inches(0.24)
        para(tf, title, size=13, color=NAVY, bold=True, first=True, space_after=4, line=1.14)
        para(tf, body, size=10.5, color=GRAY, space_after=0, line=1.22)

    stats = [("3", "agent harnesses normalized", BLUE),
             ("12", "language-model tools + 1 chat skill", PURPLE),
             ("100%", "local \u2014 nothing leaves the machine", TEAL)]
    sx = MARGIN
    sw = Inches(3.83)
    for val, label, col in stats:
        c = card(s, sx, Inches(5.54), sw, Inches(1.1), fill=WHITE, line=LINE)
        rect(s, sx, Inches(5.54), sw, Inches(0.055), fill=col)
        tf = c.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.22)
        para(tf, val, size=26, color=col, font=FONT_LT, bold=True, first=True, space_after=1)
        para(tf, label, size=11, color=GRAY, space_after=0, line=1.15)
        sx += sw + Inches(0.2)

    footer(s, FOOT)
    notes(s, """
SCRIPT

Claims only. Do NOT re-explain anything \u2014 they just watched it work. Read three of the six.

"So what's true now that wasn't in June."

"Agent sessions are a first-class, queryable object. On your own machine, across three
harnesses, with one definition of a turn, a model call, a tool call and a token."

"Semantic failures are diagnosable. 'It misunderstood me on turn five' went from a feeling to
something you can point at."

"And the agent can reason about itself. 'Why was that slow', 'compare these two runs' \u2014 those
are answerable inside the editor now, and the answers are clickable."

[say this one deliberately \u2014 it's a design position, not a limitation]
"And all of it is local by construction. Nothing leaves your machine \u2014 which matters precisely
because the interesting data here is your prompts, the model's replies, your tool arguments and
your source code. That's exactly the data you cannot ship to a shared collector."

[the three numbers, one breath \u2014 don't labour them]
"Three harnesses normalized. Twelve tools plus a skill. And none of it leaves your machine."

"Those are the three facts worth keeping if you forget everything else."

Then: "and here's where it goes from here."

\u2500\u2500\u2500 DETAIL FOR Q&A \u2500\u2500\u2500

\u2022 Content capture \u2014 prompts and replies \u2014 is opt-in, and there's a threat model in the repo.
\u2022 "It survives real use" covers multi-window handover, bounded retention with a per-service
  floor, and the worker thread. There's an appendix slide on all of it.
\u2022 If someone asks for the full inventory of what shipped \u2014 packages, explorers, tests, the
  .vsix \u2014 that's an appendix slide, and it's a good one to jump to.
\u2022 Retention is per-service with a floor deliberately: global pruning would evict Claude Code
  before Copilot purely because Copilot is chattier, which would silently bias exactly the
  cross-agent comparison the product exists to provide.
""")
    return s


def s_scorecard(prs):
    s = blank(prs)
    y = header(s, "Appendix \u00b7 scope", "Everything that made it into the box",
               "One extension, four packages, two consumer surfaces \u2014 installable today.")
    rows = [
        ["Capability", "Form it took", "Where it lives"],
        ["Local-first OTel collection", "An OTLP receiver embedded in the extension \u2014 "
         "nothing leaves the machine", "packages/receiver"],
        ["Trace, log and metric explorers", "Three explorers over a local SQLite store, "
         "on a worker thread", "Traces / Logs / Metrics tabs"],
        ["A session model across three harnesses", "One definition of session, turn, model call, "
         "tool call and token", "packages/engine"],
        ["Sessions and Home", "Two surfaces the brief never asked for \u2014 and the ones people "
         "actually open", "Sessions / Home tabs"],
        ["An agent-facing surface", "12 vscode.lm tools plus a bundled skill, every answer "
         "carrying a deeplink back into the panel", "packages/extension"],
        ["A live cost signal", "Status bar showing today\u2019s tokens and cache hit rate, "
         "refreshed as runs land", "Status bar"],
        ["Installable build", "agent-otel-0.2.0.vsix, rebuilt on every push to main",
         "CI artifact"],
        ["Tested against all three harnesses", "522 assertions across 15 suites, with fixtures "
         "captured per harness", "test/"],
    ]
    tbl = add_table(s, rows, MARGIN, y + Inches(0.02), Inches(11.89),
                    col_w=[3.5, 6.1, 2.29], row_h=Inches(0.44), header_h=Inches(0.4))
    style_table(tbl, size=10.5, header_size=11)
    for r in range(1, len(rows)):
        cell = tbl.cell(r, 2)
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.color.rgb = BLUE
                run.font.name = MONO
                run.font.size = Pt(9.5)

    footer(s, FOOT + "  \u00b7  appendix")
    notes(s, """
APPENDIX \u2014 not in the running order. Jump here if someone asks what "an extension" actually
means in scope terms, or wants the inventory of what was built.

Do NOT read this table. Say three things and move:

  "This is the whole thing in one view. Collection, storage, three explorers, the session model,
   two surfaces on top of it, and a build anyone can install."

  Point at row 4 (Sessions and Home): "these two weren't in the original brief at all. They came
   out of using it. They're now the two tabs people actually open."

  Point at the last row: "and all of it is tested against fixtures captured from all three
   harnesses at once, so the conventions can't quietly drift apart."
""")
    return s


def s_next(prs):
    s = blank(prs)
    y = header(s, "05  Impact \u00b7 what\u2019s next", "Where this goes from here",
               "Ranked by what I\u2019d pick up first.")
    items = [
        ("Codex assistant text", MAGENTA,
         "Codex never exports what the model said. The one gap nothing downstream can close \u2014 "
         "and an upstream ask worth making."),
        ("Cross-process OTel inside VS Code", PURPLE,
         "The natural next surface \u2014 and cheaper now, because this extension is already the "
         "consumer it needs."),
        ("Cost, not just tokens", AMBER,
         "Per-model pricing turns the status bar from a usage baseline into a spend baseline."),
        ("Session diffing as a first-class view", BLUE,
         "\u201cSame prompt, two agents\u201d \u2014 currently something you assemble by reading two summaries "
         "side by side."),
        ("Longer-horizon history", TEAL,
         "The projection pattern already survives pruning. It gives weeks of trend data without "
         "weeks of spans."),
        ("Discovery, and a Marketplace listing", GREEN,
         "Install tips at the moments that matter. The build already ships on every push."),
    ]
    cw = Inches(5.83)
    for i, (title, col, body) in enumerate(items):
        x = MARGIN + (cw + Inches(0.23)) * (i % 2)
        yy = y + Inches(0.02) + Inches(1.52) * (i // 2)
        c = card(s, x, yy, cw, Inches(1.34), fill=WHITE, line=LINE)
        rect(s, x, yy, cw, Inches(0.055), fill=col)
        nb = card(s, x + Inches(0.26), yy + Inches(0.44), Inches(0.4), Inches(0.4),
                  fill=col, line=col, radius=0.5)
        fill_shape_text(nb, [(str(i + 1), dict(size=13, color=WHITE, bold=True,
                                               align=PP_ALIGN.CENTER, space_after=0))])
        tf = textbox(s, x + Inches(0.84), yy + Inches(0.3), cw - Inches(1.1), Inches(1.0))
        para(tf, title, size=12.5, color=NAVY, bold=True, first=True, space_after=4, line=1.14)
        para(tf, body, size=10.5, color=GRAY, space_after=0, line=1.2)

    footer(s, FOOT)
    notes(s, """
SCRIPT

Fast \u2014 twenty seconds. This is a roadmap, not an apology. Deliver it that way.

"Six things I'd pick up next, and I'll name two."

[#1 \u2014 this is a genuine upstream ask, and you want someone in the room to hear it]
"The first one isn't something I can fix. Codex strips the model's own words before export \u2014 it
never sends them at all. Everything else in this project works around vendor differences. That
one is a hole nothing downstream can fill, and it's worth raising upstream."

[#2 \u2014 frame as a head start you're handing over, NOT as something left undone]
"The second is cross-process telemetry inside VS Code itself. That's the natural next surface,
and it's a lot cheaper to build now than it was in June \u2014 because the thing that consumes it
already exists."

"The rest are on the slide."

\u2500\u2500\u2500 DETAIL FOR Q&A \u2500\u2500\u2500

\u2022 If asked "what would you do with another month?" \u2014 #2 and #3. Cross-process OTel, and per-model
  pricing so the status bar shows spend rather than usage.
\u2022 #4, session diffing: "same prompt, two agents" or "before and after my prompt change". The data
  is already there; it's a view, not an engine change.
\u2022 #5: the projection tables already survive retention pruning, which is what makes long-horizon
  trends possible without storing months of raw spans.
\u2022 #6: the .vsix builds on every push to main. Listing it on the Marketplace is genuinely the
  last mile, not a rewrite.
""")
    return s


def s_learn1(prs):
    s = blank(prs)
    y = header(s, "Appendix \u00b7 learnings", "The longer version, with the reasoning",
               "Not in the running order \u2014 here for Q&A.")
    items = [
        ("1", "Build for what the data shows, not for who sent it", BLUE,
         "Every filter written as \u201cif it\u2019s Codex, do X\u201d would have rotted within a release. "
         "ECHO_TRACE is defined by shape \u2014 has content, no prompt, no round trip, no model-call "
         "span \u2014 and that is why it still holds. Harness-specific code is a dated cheque."),
        ("2", "A platform moving under you changes which problem you\u2019re solving", MAGENTA,
         "The agent host arrived mid-project and made collection easier and interpretation harder. "
         "The lesson isn\u2019t \u201cplatforms are unstable\u201d \u2014 it\u2019s that the moment the plumbing gets "
         "solved for you, the remaining value moves up a layer, into semantics. Plan to be there."),
        ("3", "The hard part was never the pipeline", PURPLE,
         "The receiver is ~200 lines. The normalization layer is ~2 600. That ratio is the "
         "finding: OTel collection is commodity; making three vendors\u2019 emissions mean the same "
         "thing is not."),
        ("4", "Choosing what not to build is part of building it", AMBER,
         "I aimed the whole project at agent sessions rather than generic app traces, and went "
         "deep on one surface instead of wide on several. Scope that narrowed on one axis went "
         "much further on another \u2014 and I can point at the reason."),
    ]
    h = Inches(1.06)
    for i, (num, title, col, body) in enumerate(items):
        yy = y + Inches(0.02) + (h + Inches(0.13)) * i
        c = card(s, MARGIN, yy, Inches(11.89), h, fill=WHITE, line=LINE)
        rect(s, MARGIN, yy, Inches(0.07), h, fill=col)
        nb = card(s, MARGIN + Inches(0.3), yy + Inches(0.32), Inches(0.48), Inches(0.48),
                  fill=col, line=col, radius=0.5)
        fill_shape_text(nb, [(num, dict(size=15, color=WHITE, bold=True,
                                        align=PP_ALIGN.CENTER, space_after=0))])
        tf = textbox(s, MARGIN + Inches(1.0), yy + Inches(0.2), Inches(10.6), Inches(0.9))
        para(tf, title, size=14.5, color=NAVY, bold=True, first=True, space_after=5)
        para(tf, body, size=11.5, color=GRAY, space_after=0, line=1.22)

    footer(s, FOOT + "  \u00b7  appendix")
    notes(s, """
APPENDIX \u2014 not in the running order. The four lessons in full, with the reasoning, for Q&A.

The condensed version of #1, #3 and #4 is on the "What I learned" slide. Come here if someone
asks a follow-up, or if you want #2 \u2014 the platform lesson \u2014 which isn't on the main slide.

#2 is worth having ready: the agent host arrived mid-project and made collection easier and
interpretation harder. The lesson isn't "platforms are unstable" \u2014 it's that the moment the
plumbing gets solved for you, the remaining value moves up a layer into semantics. Plan to
be there.
""")
    return s


def s_learn2(prs):
    s = blank(prs)
    y = header(s, "Appendix \u00b7 learnings", "Four more, all of them things that bit me",
               "Not in the running order \u2014 here for Q&A.")
    items = [
        ("5", "Measure before optimizing \u2014 then measure the optimizer", TEAL,
         "ANALYZE was an 8.4 s \u2192 0.2 s difference and it is one statement. Column order in a "
         "SQLite table was an order of magnitude on scans. Neither was where I would have guessed."),
        ("6", "Blocking the extension host is a product bug, not a perf bug", RED,
         "A slow query froze the whole editor \u2014 not the panel, the editor. Moving the database "
         "onto a worker thread was the single change that made this feel like a real extension."),
        ("7", "Silence is the worst failure mode", AMBER,
         "A port mismatch dropped every span with no error anywhere. A thrown query error vanished "
         "into a console the user never sees and left \u201cloading spans\u2026\u201d on screen forever. Both "
         "fixes were about making failure visible, and both mattered more than any feature."),
        ("8", "Design for the agent as a first-class consumer", GREEN,
         "Writing tool output for an LLM is its own discipline: page and budget instead of "
         "truncating; say \u201ccontent capture was off \u2014 do not infer what was said\u201d instead of "
         "returning an empty list; make deeplinks mandatory so an answer stays attached to the "
         "thing it describes."),
    ]
    h = Inches(1.06)
    for i, (num, title, col, body) in enumerate(items):
        yy = y + Inches(0.02) + (h + Inches(0.13)) * i
        c = card(s, MARGIN, yy, Inches(11.89), h, fill=WHITE, line=LINE)
        rect(s, MARGIN, yy, Inches(0.07), h, fill=col)
        nb = card(s, MARGIN + Inches(0.3), yy + Inches(0.32), Inches(0.48), Inches(0.48),
                  fill=col, line=col, radius=0.5)
        fill_shape_text(nb, [(num, dict(size=15, color=WHITE, bold=True,
                                        align=PP_ALIGN.CENTER, space_after=0))])
        tf = textbox(s, MARGIN + Inches(1.0), yy + Inches(0.2), Inches(10.6), Inches(0.9))
        para(tf, title, size=14.5, color=NAVY, bold=True, first=True, space_after=5)
        para(tf, body, size=11.5, color=GRAY, space_after=0, line=1.22)

    footer(s, FOOT + "  \u00b7  appendix")
    notes(s, """
[Appendix \u2014 not in the running order. Extra learnings #6 and #7, kept for Q&A.]

Land #6 and #7.

#6: the bug report was "VS Code freezes". Not "the panel is slow" \u2014 the whole editor, because
sql.js was running on the extension host thread and a multi-second query blocked it. That
reframed it from a perf ticket into a product bug, and the worker thread was the fix.

#7 is the one I'd generalise furthest beyond this project. Two separate failures, same shape:
the system was broken and told nobody. A port mismatch silently dropped 100% of telemetry. A
thrown query error disappeared into a console the user never opens, leaving "loading spans..."
forever. Neither was a hard bug to fix; both were invisible until someone complained. Now the
extension detects the port mismatch and OFFERS to fix it \u2014 it doesn't edit your settings behind
your back \u2014 and every webview request is answered or turned into a visible error.
""")
    return s


def s_ack(prs):
    s = blank(prs)
    bg(s, NAVY)
    accent_bar(s)
    tf = textbox(s, MARGIN, Inches(0.55), CONTENT_W, Inches(0.3))
    para(tf, "ACKNOWLEDGEMENTS", size=10.5, color=BLUE_L, bold=True, font=FONT_SB, first=True,
         space_after=0)
    tf = textbox(s, MARGIN, Inches(0.9), CONTENT_W, Inches(0.6))
    para(tf, "Standing on other people\u2019s shoulders, specifically",
         size=29, color=WHITE, font=FONT_LT, first=True, space_after=0)
    rect(s, MARGIN, Inches(1.62), Inches(1.05), Inches(0.035), fill=TEAL)

    people = [
        ("@zhichli", "Author of the original OTel Studio for VS Code proposal and the otelux "
                     "prototype that framed the problem space \u2014 and co-assignee on the issue.", BLUE_L),
        ("@karthiknadig  \u00b7  @kieferrm", "For the early direction on the proposal.", TEAL),
    ]
    yy = Inches(2.0)
    for name, body, col in people:
        c = card(s, MARGIN, yy, Inches(11.89), Inches(0.86),
                 fill=RGBColor(0x1E, 0x33, 0x54), line=RGBColor(0x30, 0x45, 0x6B))
        rect(s, MARGIN, yy, Inches(0.06), Inches(0.86), fill=col)
        tf = c.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.3)
        para(tf, [(name, dict(size=14, color=col, bold=True)),
                  ("     " + body, dict(size=12, color=RGBColor(0xC3, 0xD1, 0xE6)))],
             first=True, space_after=0, line=1.22)
        yy += Inches(0.98)

    teams = [
        ("The VS Code agent-host team", "For microsoft/vscode#328529 \u2014 the session anchor span \u2014 "
         "and the chat.agentHost.otel.* surface this is built on. And for #316090, which is what "
         "makes Stream 2 a real next step rather than an idea.", PURPLE),
        ("The OpenTelemetry GenAI semconv WG", "gen_ai.* is the only reason a shared vocabulary "
         "exists to normalize toward. Without it this project is three parsers and a prayer.", MAGENTA),
        ("The Copilot, Claude Code and Codex teams", "For emitting telemetry at all. Every "
         "difference in that comparison table exists because three teams each independently "
         "decided this was worth instrumenting.", AMBER),
    ]
    cw = Inches(3.83)
    for i, (name, body, col) in enumerate(teams):
        x = MARGIN + (cw + Inches(0.2)) * i
        c = card(s, x, yy, cw, Inches(1.62), fill=RGBColor(0x1E, 0x33, 0x54),
                 line=RGBColor(0x30, 0x45, 0x6B))
        rect(s, x, yy, cw, Inches(0.055), fill=col)
        tf = c.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.24)
        para(tf, name, size=13, color=col, bold=True, first=True, space_after=6, line=1.15)
        para(tf, body, size=11, color=RGBColor(0xB9, 0xC8, 0xE0), space_after=0, line=1.22)

    tf = textbox(s, MARGIN, Inches(6.42), Inches(11.89), Inches(0.4))
    para(tf, "\u2014 and everyone who dogfooded it and reported that their session list had "
             "261 rows in it.",
         size=13, color=RGBColor(0x8E, 0xA4, 0xC6), italic=True, first=True, space_after=0,
         align=PP_ALIGN.CENTER)

    notes(s, """
Short, specific, sincere. Do not rush it and do not read the small text.

Name Zhichao first and properly \u2014 the proposal and the otelux prototype are what framed this
problem space, and a lot of what I built is an argument with ideas he'd already had.

Then the agent-host team: their anchor span is genuinely the thing that made cross-harness
session identity tractable at all, even though \u2014 as I said earlier \u2014 it also created four new
problems for me. Both of those are true at once and it's worth saying both.

The semconv working group line is the one I mean most: gen_ai.* is the only reason there is
anything to normalize TOWARD. Without a shared target vocabulary this is three parsers and a
prayer.

End with the 261 line \u2014 it lands as a laugh and it hands the room back to you for Q&A.
""")
    return s


def s_close(prs):
    s = blank(prs)
    bg(s, NAVY)
    rect(s, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill=BLUE)
    xs = [10.05, 10.62, 11.19, 11.76, 12.33]
    hs = [1.2, 2.1, 0.8, 1.7, 1.05]
    cols = [BLUE, TEAL, PURPLE, BLUE_L, MAGENTA]
    for x, h, c in zip(xs, hs, cols):
        rect(s, Inches(x), Inches(5.9 - h), Inches(0.30), Inches(h), fill=c,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.35)

    tf = textbox(s, Inches(1.05), Inches(1.5), Inches(8.6), Inches(0.3))
    para(tf, "THANK YOU", size=12, color=BLUE_L, bold=True, font=FONT_SB, first=True,
         space_after=0)
    tf = textbox(s, Inches(1.0), Inches(1.95), Inches(8.8), Inches(2.4))
    para(tf, [("The data layer is solved.", dict(color=WHITE))],
         size=40, font=FONT_LT, first=True, space_after=6, line=1.1)
    para(tf, [("Agent Insights is about ", dict(color=RGBColor(0x93, 0xA6, 0xC4))),
              ("everything after.", dict(color=BLUE_L, bold=True))],
         size=40, font=FONT_LT, space_after=0, line=1.1)

    rect(s, Inches(1.05), Inches(4.32), Inches(1.5), Inches(0.045), fill=TEAL)

    tf = textbox(s, Inches(1.05), Inches(4.66), Inches(9.0), Inches(1.4))
    para(tf, "One definition of a session, a turn, a model call, a tool call and a token \u2014 "
             "across three harnesses, on your own machine, readable by you and by your agent.",
         size=15, color=RGBColor(0xB9, 0xC8, 0xE0), first=True, space_after=14, line=1.28)
    para(tf, [("github.com/michiisai/agent-insights", dict(font=MONO, color=BLUE_L)),
              ("      \u00b7      ", dict(color=RGBColor(0x60, 0x76, 0x9C))),
              ("code --install-extension michiisai.agent-otel", dict(font=MONO, color=RGBColor(0x93, 0xA6, 0xC4)))],
         size=12, space_after=0)

    tf = textbox(s, Inches(1.05), Inches(6.32), Inches(9.0), Inches(0.4))
    para(tf, "Questions \u2014 and there\u2019s an appendix with the numbers if you want detail.",
         size=12, color=RGBColor(0x60, 0x76, 0x9C), italic=True, first=True, space_after=0)

    notes(s, """
Close, then stop talking.

Deliver the two lines and let the second one sit for a beat:
"The data layer is solved. Agent Insights is about everything after."

That's the same sentence I opened with on slide 2, which is deliberate \u2014 it's the frame for the
whole talk.

Then: "Happy to take questions. There's an appendix with all the numbers if anyone wants to go
deeper on the normalization."

LIKELY QUESTIONS \u2014 have these ready:
\u2022 "Why not MCP?" \u2014 vscode.lm tools are first-party: no separate server process, they appear as
  #-references, and the bundled skill lets the agent route without being named. MCP was the
  proposal's plan; LM tools are strictly better inside VS Code.
\u2022 "Why sql.js and not better-sqlite3?" \u2014 no native module to build per platform/Electron ABI.
  Cost is a WASM boundary, which is why it lives on a worker thread.
\u2022 "How do you know your numbers are right?" \u2014 522 assertions across 15 suites, including one
  that asserts Copilot, Claude and Codex token totals simultaneously from one fixture.
\u2022 "Privacy?" \u2014 loopback-only, content capture opt-in, threat model in docs/threat-model.md.
\u2022 "What would you do with another month?" \u2014 Stream 2, and per-model pricing.
""")
    return s


# ------------------------------------------------------------------ appendix
def s_appendix_numbers(prs):
    s = blank(prs)
    y = header(s, "Appendix", "Numbers worth having in your pocket")
    rows = [
        ["Metric", "Value", "Metric", "Value"],
        ["Phantom vs real sessions, one day", "261  vs  6", "LM tools / chat skills", "12  /  1"],
        ["Codex spans per single model call", "5 (nested)", "LM tool timeout \u2014 always settles", "15 s"],
        ["ANALYZE on the recursive trace walk", "8.4 s \u2192 0.2 s", "Transcript budget", "10 turns / 1 500 ch / 40 000 cap"],
        ["idx_raw_spans_trace, 3k-span trace", "~10 s \u2192 fast", "Collector heartbeat / takeover", "2 s \u00b7 2 fails \u00b7 \u2264750 ms jitter"],
        ["Retention cap", "50 000 rows \u00b7 96 MB spans", "Default port", "4318"],
        ["Per-service row floor", "5 000", "Normalization vs receiver", "~2 600 vs ~200 lines"],
        ["token_facts retention", "9 days (survives pruning)", "Source / test lines", "17 372  /  3 696"],
        ["Commits \u00b7 span of work", "239 \u00b7 30 Jun \u2013 19 Aug", "Test assertions passing", "522 across 15 suites"],
    ]
    tbl = add_table(s, rows, MARGIN, y + Inches(0.02), Inches(11.89),
                    col_w=[3.4, 2.6, 3.4, 2.6], row_h=Inches(0.42), header_h=Inches(0.42))
    style_table(tbl, size=11, header_size=11)
    for r in range(1, len(rows)):
        for c in (1, 3):
            for p in tbl.cell(r, c).text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = MONO
                    run.font.bold = True
                    run.font.color.rgb = BLUE
        for p in tbl.cell(r, 2).text_frame.paragraphs:
            for run in p.runs:
                run.font.bold = True
                run.font.color.rgb = NAVY
    footer(s, FOOT + "  \u00b7  appendix")
    notes(s, """
Backup slide. Jump here for "how big was it / how do you know it's right" questions.

Favourites to quote: 261 vs 6. 8.4 s \u2192 0.2 s from one ANALYZE. And ~2 600 lines of normalization
against ~200 lines of receiver \u2014 the ratio that is the whole finding.
""")
    return s


def s_appendix_setup(prs):
    s = blank(prs)
    y = header(s, "Appendix", "Setup \u2014 the whole configuration",
               "Two settings blocks, one reload, and every session after that is captured.")
    grey = RGBColor(0x7A, 0x93, 0xB8)
    code_block(s, MARGIN, y + Inches(0.02), Inches(7.1), Inches(3.5), [
        [("{", {})],
        [("  // Native Agent Host: Copilot, Claude Code, Codex", dict(color=grey))],
        [('  "chat.agentHost.enabled"', dict(color=BLUE_L)), (': true,', {})],
        [('  "chat.agentHost.otel.enabled"', dict(color=BLUE_L)), (': true,', {})],
        [('  "chat.agentHost.otel.captureContent"', dict(color=BLUE_L)), (': true,', {})],
        [('  "chat.agentHost.otel.otlpEndpoint"', dict(color=BLUE_L)),
         (': "http://localhost:4318",', dict(color=AMBER))],
        "",
        [("  // Extension-side Copilot / VS Code LM telemetry", dict(color=grey))],
        [('  "github.copilot.chat.otel.enabled"', dict(color=BLUE_L)), (': true,', {})],
        [('  "github.copilot.chat.otel.captureContent"', dict(color=BLUE_L)), (': true,', {})],
        [('  "github.copilot.chat.otel.otlpEndpoint"', dict(color=BLUE_L)),
         (': "http://localhost:4318"', dict(color=AMBER))],
        [("}", {})],
    ], size=12)

    rx = MARGIN + Inches(7.36)
    cards_ = [
        ("Content capture is opt-in, and says so", MAGENTA,
         "It may include prompts, responses, tool arguments and file contents. Leave it off when "
         "exporting to a shared or untrusted collector. Documented in docs/threat-model.md."),
        ("Claude and Codex need no separate config", GREEN,
         "Once the agent host is on, they inherit the host\u2019s exporter. Only Copilot\u2019s "
         "extension-side telemetry has its own block."),
        ("Two settings, two signal sets", AMBER,
         "Copilot emits traces from one setting and metrics/logs from another \u2014 so one can work "
         "while the other is off. The commonest \u201cno data\u201d cause is a port mismatch, which the "
         "extension detects and offers to fix rather than silently editing your settings."),
    ]
    yy = y + Inches(0.02)
    for title, col, body in cards_:
        c = card(s, rx, yy, Inches(4.53), Inches(1.1), fill=WHITE, line=LINE)
        rect(s, rx, yy, Inches(0.055), Inches(1.1), fill=col)
        tf = c.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.26)
        tf.margin_right = Inches(0.22)
        para(tf, title, size=12, color=NAVY, bold=True, first=True, space_after=4, line=1.14)
        para(tf, body, size=10, color=GRAY, space_after=0, line=1.2)
        yy += Inches(1.2)

    footer(s, FOOT + "  \u00b7  appendix")
    notes(s, """
Backup slide, and also the "how do I try it" slide if anyone asks after the talk.

The point to make if it comes up: this is the entire setup. Two blocks, one reload. There is no
collector to run, no Docker, no account. That low a setup cost is the reason it gets used at all
\u2014 the alternative was "stand up Jaeger", and nobody does that to look at their own chat history.
""")
    return s


def s_appendix_hard(prs):
    s = blank(prs)
    y = header(s, "Appendix", "The five hardest normalization problems, and where they\u2019re solved")
    rows = [
        ["Problem", "Mechanism", "Where"],
        ["Codex traces rejoin their session \u2014 the host anchors only the first trace of a thread",
         "codex_trace_sessions: seed a session alias per trace from Codex\u2019s own conversation.id log "
         "records, then promote the whole conversation to the host\u2019s id once any trace turns out to "
         "be anchored", "store.ts  \u00b7  sessions.ts"],
        ["Codex echo-trace duplication \u2014 each host-executed tool call is logged a second time on "
         "its own trace",
         "ECHO_TRACE structural predicate: content events, but no prompt, no round trip and no "
         "model-call span. Excluded from sessions, still browsable in Traces", "sessions.ts"],
        ["Claude call-order reconstruction \u2014 api_request is logged when the round trip finishes, "
         "so records arrive before the call they close",
         "ClaudeCallDraft adopts an already-open unclaimed draft instead of minting a new call, so "
         "a 7-call agent loop splits into 7 turns rather than collapsing into 1", "sessions.ts"],
        ["Three cache-token conventions at once, plus parent spans that re-report their children",
         "is_additive flag per span in token_facts + conditional expansion at query time; rollup "
         "spans excluded at harvest by gen_ai.operation.name", "store.ts  \u00b7  tokenRows.ts"],
        ["Agent-host span inflation and multi-turn projection \u2014 one physical trace holds many turns",
         "AGENT_SPAN_COUNT and AGENT_SERVICE_NAME exclude host spans; getTraces() promotes each "
         "direct child of a host root into a logical trace <traceId>:<rootSpanId>",
         "sessions.ts  \u00b7  traces.ts"],
    ]
    tbl = add_table(s, rows, MARGIN, y + Inches(0.02), Inches(11.89),
                    col_w=[3.6, 6.3, 2.0], row_h=Inches(0.92), header_h=Inches(0.4))
    style_table(tbl, size=10, header_size=11)
    for r in range(1, len(rows)):
        for p in tbl.cell(r, 2).text_frame.paragraphs:
            for run in p.runs:
                run.font.name = MONO
                run.font.size = Pt(9.5)
                run.font.color.rgb = BLUE
    footer(s, FOOT + "  \u00b7  appendix")
    notes(s, """
Backup slide for a technical audience \u2014 use it if someone says "give me a concrete example of a
normalization problem."

Best single story to tell from here is row 3, the Claude one, because it's counter-intuitive:
Claude logs the api_request record when the round trip FINISHES, not when it starts. So a
permission decision or a tool result can legitimately show up in the stream before the request
that owns it. If you naively open a new call on every api_request, a seven-call agent loop
collapses into one turn with all seven tools piled onto the last call. The fix is that opening a
call adopts an already-open, unclaimed draft.

Every one of these five has a dedicated test suite. That's most of the 522 assertions.
""")
    return s
