"""Slides: architecture, storage, surfaces, normalization."""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from theme import *



# ------------------------------------------------------------------ architecture
def s_architecture(prs):
    s = blank(prs)
    y = header(s, "03  Approach \u00b7 architecture", "Four packages, one pipeline, two consumers",
               "The split is what keeps the SQL out of the UI and the UI out of the agent.")

    L, W = MARGIN, Inches(11.89)
    y = Inches(1.74)
    # row 1 - producers
    sources = [("GitHub Copilot", "github-copilot", BLUE),
               ("Claude Code", "claude-code", AMBER),
               ("Codex", "codex-app-server", MAGENTA),
               ("VS Code Agent Host", "vscode-agent-host", PURPLE)]
    bw = Inches(2.86)
    gap = Inches(0.144)
    for i, (name, svc, col) in enumerate(sources):
        x = L + (bw + gap) * i
        sh = card(s, x, y, bw, Inches(0.58), fill=WHITE, line=col, radius=0.07)
        sh.line.width = Pt(1.5)
        tf = sh.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.12)
        para(tf, name, size=12, color=NAVY, bold=True, align=PP_ALIGN.CENTER, first=True,
             space_after=1)
        para(tf, svc, size=9, color=col, font=MONO, align=PP_ALIGN.CENTER, space_after=0)

    y2 = y + Inches(0.58)
    down_arrow(s, L + W / 2, y2 + Inches(0.05), color=RGBColor(0xB3, 0xB0, 0xAD),
               w=Inches(0.18), h=Inches(0.24))
    tf = textbox(s, L, y2 + Inches(0.08), Inches(5.6), Inches(0.24))
    para(tf, "OTLP / HTTP JSON", size=10, color=GRAY, font=MONO, bold=True,
         align=PP_ALIGN.RIGHT, first=True, space_after=0)
    tf = textbox(s, L + W / 2 + Inches(0.3), y2 + Inches(0.08), Inches(5.6), Inches(0.24))
    para(tf, "127.0.0.1:4318   \u00b7   loopback only", size=10, color=GRAY, font=MONO,
         first=True, space_after=0)

    def band(yy, h, title, subtitle, col, fill=WHITE, tcol=None, scol=None, tw=None):
        sh = card(s, L, yy, W, h, fill=fill, line=LINE, radius=0.05)
        rect(s, L, yy, Inches(0.075), h, fill=col)
        tf = sh.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.34)
        tf.margin_right = tw or Inches(0.3)
        para(tf, title, size=13, color=tcol or NAVY, bold=True, first=True, space_after=3)
        para(tf, subtitle, size=10, color=scol or GRAY, font=MONO, space_after=0, line=1.2)
        return sh

    yy = y2 + Inches(0.4)
    band(yy, Inches(0.7), "@agent-insights/receiver",
         "/v1/traces  \u00b7  /v1/metrics  \u00b7  /v1/logs  \u00b7  /.well-known/agent-insights   "
         "\u2014  flattens the OTLP envelope to one row per entity, preserving values exactly",
         BLUE)

    yy2 = yy + Inches(0.7)
    down_arrow(s, L + W / 2, yy2 + Inches(0.03), color=RGBColor(0xB3, 0xB0, 0xAD),
               w=Inches(0.18), h=Inches(0.24))
    tf = textbox(s, L + W / 2 + Inches(0.26), yy2 + Inches(0.06), Inches(5.4), Inches(0.24))
    para(tf, "worker_threads boundary \u2014 the extension host never touches SQL", size=10,
         color=RED, font=MONO, first=True, space_after=0)

    yy = yy2 + Inches(0.36)
    band(yy, Inches(0.74), "sql.js  \u00b7  SQLite compiled to WASM, on a worker thread",
         "raw OTLP JSON per row + materialized derived columns   \u00b7   projections that outlive "
         "retention:  session_titles  \u00b7  codex_trace_sessions  \u00b7  token_facts", TEAL)

    yy2 = yy + Inches(0.74)
    down_arrow(s, L + W / 2, yy2 + Inches(0.03), color=RGBColor(0xB3, 0xB0, 0xAD),
               w=Inches(0.18), h=Inches(0.24))

    yy = yy2 + Inches(0.34)
    sh = band(yy, Inches(0.78), "@agent-insights/engine   \u2014   the normalization layer",
              "sessions  \u00b7  traces  \u00b7  tokenUsage  \u00b7  agentAnalytics  \u00b7  metrics  \u00b7  logs",
              PURPLE, fill=NAVY, tcol=WHITE, scol=RGBColor(0x9E, 0xB4, 0xD6),
              tw=Inches(1.9))
    ch = card(s, Inches(10.78), yy + Inches(0.19), Inches(1.6), Inches(0.4),
              fill=PURPLE, line=PURPLE, radius=0.5)
    fill_shape_text(ch, [("\u2605 the contribution", dict(size=10, color=WHITE, bold=True,
                                                          align=PP_ALIGN.CENTER, space_after=0))])

    yy2 = yy + Inches(0.78)
    down_arrow(s, L + Inches(2.9), yy2 + Inches(0.03), color=RGBColor(0xB3, 0xB0, 0xAD),
               w=Inches(0.18), h=Inches(0.24))
    down_arrow(s, L + Inches(8.95), yy2 + Inches(0.03), color=RGBColor(0xB3, 0xB0, 0xAD),
               w=Inches(0.18), h=Inches(0.24))

    yy = yy2 + Inches(0.34)
    cw = Inches(5.83)
    for i, (title, sub, col, who) in enumerate([
        ("Webview panel", "Home \u00b7 Sessions \u00b7 Traces \u00b7 Metrics \u00b7 Logs", BLUE, "for humans"),
        ("12 language-model tools + 1 skill",
         "#agentSession  #agentTranscript  #agentService  \u2026", GREEN, "for agents")]):
        x = L + (cw + Inches(0.23)) * i
        sh = card(s, x, yy, cw, Inches(0.7), fill=WHITE, line=col, radius=0.06)
        sh.line.width = Pt(1.5)
        tf = sh.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.26)
        para(tf, [(title, dict(color=NAVY, bold=True, size=12.5)),
                  ("   " + who, dict(color=col, size=10, bold=True))],
             first=True, space_after=2)
        para(tf, sub, size=9.5, color=GRAY, font=MONO, space_after=0)

    tf = textbox(s, L, yy + Inches(0.76), W, Inches(0.3))
    para(tf, [("\u2194  ", dict(color=GREEN, bold=True, size=12)),
              ("both directions: ", dict(color=NAVY, bold=True)),
              ("every tool result carries a deeplink  ", {}),
              ("vscode://\u2026/navigate?sessionId=\u2026", dict(font=MONO, color=BLUE)),
              ("  \u00b7  and \u201c+ chat\u201d stages a pre-built query back into the chat input", {})],
         size=10, color=GRAY, align=PP_ALIGN.CENTER, first=True, space_after=0)

    footer(s, FOOT)
    notes(s, """
SCRIPT

The hero diagram. Trace it top to bottom with your hand, roughly ten seconds per layer.

"This is the whole system on one slide. Four packages, one pipeline, two consumers."

[top]
"At the top, four producers. The three agents, plus the host itself. All of them speak OTLP over
plain HTTP to a port on localhost. Loopback only \u2014 there is deliberately no remote mode."

[receiver]
"The receiver is a plain Node HTTP server. About two hundred lines. It flattens the OTLP envelope
to one row per entity and preserves the values exactly as sent."

[the red boundary \u2014 worth a beat]
"Then a worker-thread boundary. SQLite, compiled to WebAssembly, running off the main thread.
It used to run in-process, and a multi-second query froze the entire editor \u2014 not the panel, the
editor. That's a product bug, not a performance bug."

[engine \u2014 slow down, this is the point of the slide]
"And then the engine. Everything above this line is plumbing you can get off the shelf. This
layer is the one where three vendors' emissions are made to mean the same thing. It's why the
badge says 'the contribution'."

[bottom]
"Out the bottom, two consumers over one model \u2014 a panel for a person, tools for an agent. And
note the arrows go both ways: every tool result carries a link back into the panel, so a human
and an agent are always pointing at the same object."

Then hand off: "that engine box is the whole project. Let me show you one problem inside it
properly."

\u2500\u2500\u2500 DETAIL FOR Q&A \u2500\u2500\u2500

\u2022 Sizes: receiver ~1.4k lines, engine ~4.9k, about 17k lines total, 522 test assertions.
\u2022 Why sql.js and not better-sqlite3: no native module to rebuild per platform and per Electron
  ABI. The cost is a WASM boundary, which is exactly why it lives on a worker thread.
\u2022 Value preservation: an OTLP intValue stays a string through the parser, so a 64-bit token count
  doesn't quietly lose precision by becoming a JavaScript number.
\u2022 The four packages are receiver, engine, extension and types \u2014 the split is what keeps SQL out
  of the UI and the UI out of the agent surface.
""")
    return s


# ------------------------------------------------------------------ ingest & storage
def s_storage(prs):
    s = blank(prs)
    y = header(s, "Appendix \u00b7 engineering",
               "Ingest and storage: the part that makes it usable",
               "Nothing here is novel. All of it is why it survives a real workday.")

    cw = Inches(3.83)
    gap = Inches(0.2)
    cards = [
        ("Multi-window collector handover", BLUE,
         "A port can be owned by one window; developers run five.",
         ["Bind \u2192 on EADDRINUSE probe the occupant 3\u00d7, 100 ms apart",
          "If it\u2019s a peer, follow it read-only on a 2 s heartbeat",
          "2 failed beats \u2192 \u2264750 ms jitter \u2192 attempt takeover",
          "Outgoing owner drains: 503 + retry-after, then waitForIdle()",
          "Someone else\u2019s collector? Say so \u2014 don\u2019t fail silently"]),
        ("Raw-first schema, generated", TEAL,
         "Each row keeps one self-contained OTLP entity as JSON.",
         ["Every queryable field materialized into a real column",
          "Schema, inserts, migrations and views all generated from one DERIVED table \u2014 so they cannot drift",
          "DERIVED_VERSION per row drives re-derivation",
          "Column order is load-bearing: attributes and raw declared last, or scans walk multi-KB overflow pages first"]),
        ("Retention with a conscience", PURPLE,
         "50 000 rows / 96 MB for spans \u2014 but not globally.",
         ["Per-service pruning with a 5 000-row floor per service",
          "Otherwise Copilot\u2019s volume evicts Claude Code purely for being older \u2014 quietly biasing the agent-comparison views the product exists for",
          "Pruning refuses to orphan a referenced parent span",
          "session_titles, codex_trace_sessions and token_facts are exempt \u2014 they outlive their spans"]),
    ]
    for i, (title, col, lede, bullets) in enumerate(cards):
        x = MARGIN + (cw + gap) * i
        c = card(s, x, y, cw, Inches(3.44), fill=WHITE, line=LINE)
        rect(s, x, y, cw, Inches(0.055), fill=col)
        tf = c.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = tf.margin_right = Inches(0.24)
        tf.margin_top = Inches(0.22)
        para(tf, title, size=13.5, color=NAVY, bold=True, first=True, space_after=4, line=1.15)
        para(tf, lede, size=11, color=col, bold=True, space_after=8, line=1.2)
        for b in bullets:
            para(tf, [("\u00b7  ", dict(color=col, bold=True)), (b, dict(color=GRAY))],
                 size=10.5, space_after=5, line=1.18)

    # perf facts
    band = card(s, MARGIN, Inches(5.5), Inches(11.89), Inches(1.16), fill=NAVY, line=NAVY)
    tf = textbox(s, MARGIN + Inches(0.36), Inches(5.68), Inches(3.6), Inches(0.9))
    para(tf, "MEASURED, NOT GUESSED", size=10, color=BLUE_L, bold=True, font=FONT_SB,
         first=True, space_after=6)
    para(tf, "Neither of these was where I would have guessed.",
         size=11.5, color=RGBColor(0x93, 0xA6, 0xC4), space_after=0, line=1.2)
    facts = [("8.4 s \u2192 0.2 s", "ANALYZE. One statement. Without stats SQLite picked the wrong\nindex for the parent\u2192child hop in the recursive trace walk.", GREEN),
             ("~10 s \u2192 fast", "idx_raw_spans_trace(trace_id, parent_span_id)\non a 3 000-span trace.", BLUE_L)]
    x = Inches(4.9)
    for val, sub, col in facts:
        tf = textbox(s, x, Inches(5.7), Inches(3.6), Inches(0.9))
        para(tf, val, size=21, color=col, bold=True, font=FONT_LT, first=True, space_after=3)
        para(tf, sub.replace("\n", " "), size=10.5, color=RGBColor(0x93, 0xA6, 0xC4),
             space_after=0, line=1.2)
        x += Inches(3.75)

    footer(s, FOOT + "  \u00b7  appendix")
    notes(s, """
APPENDIX \u2014 not in the running order. Jump here if someone asks how it holds up in real use,
about multi-window behaviour, retention, or SQLite performance.

Card 1 in one sentence: "a port can only be owned by one VS Code window, and developers run
five, so the collector negotiates \u2014 bind, probe, follow read-only, heartbeat, take over cleanly
when the owner goes away. No data loss, no configuration."

Card 2, the line worth saying: schema, inserts, migrations and views are all generated from one
DERIVED table, so they physically cannot drift apart.

Card 3 is the one I'd defend in a design review: retention is per-service with a floor, because
global pruning would evict Claude Code before Copilot purely because Copilot is chattier \u2014 which
would silently bias exactly the cross-agent comparison this product exists to provide. A
correctness bug wearing a performance bug's clothes.

Bottom band: ANALYZE is a single statement worth 8.4 seconds. Column order in a SQLite table was
an order of magnitude on scans. Measure first.
""")
    return s


# ------------------------------------------------------------------ surfaces
def s_surfaces(prs):
    s = blank(prs)
    y = header(s, "03  Approach \u00b7 surfaces", "Two consumers, one model",
               "Same engine, same objects \u2014 one rendered for a person, one written "
               "for a language model.")

    # LEFT panel
    lx, lw = MARGIN, Inches(5.83)
    c = card(s, lx, y, lw, Inches(3.5), fill=WHITE, line=LINE)
    rect(s, lx, y, lw, Inches(0.055), fill=BLUE)
    tf = c.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = Inches(0.28)
    tf.margin_top = Inches(0.24)
    para(tf, "THE PANEL \u2014 FOR HUMANS", size=10.5, color=BLUE, bold=True, font=FONT_SB,
         first=True, space_after=10)
    tabs = [("Home", "totals, tokens by model, cache hit rate, errors"),
            ("Sessions", "the headline surface \u2014 list \u2192 timeline \u2192 transcript"),
            ("Traces / Metrics / Logs", "the conventional OTel views, kept")]
    for name, desc in tabs:
        para(tf, [("\u00b7  ", dict(color=BLUE, bold=True)),
                  (name + "  ", dict(color=NAVY, bold=True)), (desc, dict(color=GRAY))],
             size=11, space_after=8, line=1.2)
    para(tf, "Two product decisions, not rendering decisions", size=10.5, color=BLUE,
         bold=True, font=FONT_SB, space_after=8, space_before=8)
    para(tf, [("\u00b7  ", dict(color=BLUE, bold=True)),
              ("Shared context is hoisted ", dict(color=NAVY, bold=True)),
              ("out of turns \u2014 the harness repeats the preamble every time.",
               dict(color=GRAY))],
         size=11, space_after=7, line=1.2)
    para(tf, [("\u00b7  ", dict(color=BLUE, bold=True)),
              ("Every bubble knows its span ", dict(color=NAVY, bold=True)),
              ("\u2014 click a sentence, land on it.", dict(color=GRAY))],
         size=11, space_after=0, line=1.2)

    # RIGHT tools
    rx, rw = MARGIN + Inches(6.06), Inches(5.83)
    c = card(s, rx, y, rw, Inches(3.5), fill=WHITE, line=LINE)
    rect(s, rx, y, rw, Inches(0.055), fill=GREEN)
    tf = c.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = Inches(0.28)
    tf.margin_top = Inches(0.24)
    para(tf, "THE TOOLS \u2014 FOR AGENTS", size=10.5, color=GREEN, bold=True, font=FONT_SB,
         first=True, space_after=10)
    tools = [("#agentSession", "outcome, timeline, tokens, tools, errors"),
             ("#agentTranscript", "what was actually said, paged"),
             ("#agentService", "per-agent profile \u2014 the comparison tool"),
             ("#agentErrors  #agentSlow", "what broke, what\u2019s slow")]
    for n, d in tools:
        para(tf, [("\u00b7  ", dict(color=GREEN, bold=True)),
                  (n, dict(font=MONO, color=GREEN, bold=True, size=10.5)),
                  ("   " + d, dict(color=GRAY))],
             size=11, space_after=8, line=1.2)
    para(tf, "\u2026and eight more for traces, spans, logs and metrics.", size=11, color=GRAY,
         space_after=10, line=1.2)
    para(tf, [("A bundled skill routes the agent ", dict(color=NAVY, bold=True)),
              ("to the right tool without anyone naming it.", dict(color=GRAY))],
         size=11, space_after=0, line=1.2)

    # design choices band
    band = card(s, MARGIN, Inches(5.66), Inches(11.89), Inches(1.0), fill=WASH_2, line=LINE)
    tf = band.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.34)
    para(tf, "WRITING TOOL OUTPUT FOR AN LLM IS ITS OWN DISCIPLINE", size=10, color=GREEN,
         bold=True, font=FONT_SB, first=True, space_after=7)
    para(tf, [("Every call is guaranteed to settle.", dict(bold=True, color=NAVY)),
              ("     Paged and budgeted, never truncated.", dict(bold=True, color=NAVY)),
              ("     \u201cNo content captured\u201d is an explicit answer.", dict(bold=True, color=NAVY)),
              ("     Deeplinks are mandatory.", dict(bold=True, color=NAVY))],
         size=12, space_after=0, line=1.24)

    footer(s, FOOT)
    notes(s, """
SCRIPT

"Two consumers, one model. Everything on the left, an agent can ask for on the right \u2014 same
objects, same links."

[left]
"For a person: five tabs. Home is the daily view. Sessions is the headline surface, and it's the
one this whole talk has been about \u2014 a list, a timeline, a readable transcript. Traces, metrics
and logs are the conventional views; I kept them because sometimes you really do want the
waterfall."

"Two things here are product decisions rather than rendering decisions. The harness re-sends the
same preamble every single turn, so I hoist it out and show it once \u2014 showing it fourteen times
is noise. And every bubble in the transcript knows which span produced it, so you can click a
sentence and land on the exact operation behind it."

[right]
"For an agent: twelve tools you can hash-reference in chat. The interesting part is the skill \u2014
it routes to the right tool without the user knowing any of these names. You just ask 'why was
that slow' and it picks."

[the band]
"And writing tool output for a model turned out to be its own discipline. Four rules I'd keep on
any project: every call settles, results are paged rather than truncated, 'no content captured'
is an explicit answer, and every result carries a link back."

Then: "that's the build \u2014 let me show you it working."

\u2500\u2500\u2500 DETAIL FOR Q&A \u2500\u2500\u2500

\u2022 "Every call settles": cancellation, a 15-second timeout, and error isolation. A timeout that
  explains itself beats a promise that never resolves \u2014 a hung tool call makes the agent look
  broken.
\u2022 Paging budget: 10 turns per page, 1,500 characters each, 40,000-character ceiling. The result
  tells the model how to fetch the next page rather than silently cutting off.
\u2022 Why "no content captured" matters: return an empty list and the model infers nothing was said.
  Saying capture was off stops it inventing an answer.
\u2022 Why LM tools and not MCP: first-party, no separate server process, they appear as
  #-references, and the skill can route without the user naming a tool.
\u2022 The panel is one webview, no framework \u2014 the render surface was simple enough that a build
  step would have cost more than it saved.
""")
    return s


# ------------------------------------------------------------------ normalization I
def s_norm_identity(prs):
    s = blank(prs)
    y = header(s, "Appendix \u00b7 normalization", "Identity and counting",
               "The pattern every time: define the concept once, in SQL, from what the data shows "
               "\u2014 not from which harness sent it.")

    tf = textbox(s, MARGIN, y, Inches(6.4), Inches(0.3))
    para(tf, "1  \u00b7  WHAT IS ONE SESSION?   A COALESCE LADDER, MOST-TRUSTED FIRST",
         size=10.5, color=BLUE, bold=True, font=FONT_SB, first=True, space_after=0)

    grey = RGBColor(0x7A, 0x93, 0xB8)
    code_block(s, MARGIN, y + Inches(0.34), Inches(6.4), Inches(2.16), [
        [("COALESCE", dict(color=BLUE_L, bold=True)), ("(", {})],
        [("  MAX(json_extract(attrs,'$.\"gen_ai.conversation.id\"')),", {}),
         ("   OTel semconv", dict(color=grey))],
        [("  MAX(json_extract(attrs,'$.\"session.id\"')),", {}),
         ("               Claude", dict(color=grey))],
        [("  MAX(json_extract(attrs,'$.\"copilot_chat.chat_session_id\"')),", {})],
        [("                                                  ", {}), ("Copilot", dict(color=grey))],
        [("  (SELECT session_id FROM codex_trace_sessions ...),", {}),
         ("  Codex", dict(color=grey))],
        [("  trace_id", dict(color=AMBER)), ("                                     last resort", dict(color=grey))],
        [(")", {})],
    ], size=11)

    tf = textbox(s, MARGIN, y + Inches(2.62), Inches(6.4), Inches(1.2))
    para(tf, [("MAX(\u2026) rather than a plain read", dict(bold=True, color=NAVY)),
              (", because the id is on chat spans and absent on the permission and "
               "execute_tool spans in the same trace. One span knowing the answer is enough.",
               dict(color=GRAY))],
         size=11, first=True, space_after=7, line=1.22)
    para(tf, [("Codex often has no id at all", dict(bold=True, color=MAGENTA)),
              (" \u2014 a projection recovers one per trace from Codex\u2019s own conversation.id logs, "
               "then promotes the conversation once any trace proves anchored.",
               dict(color=GRAY))],
         size=11, space_after=0, line=1.22)

    # right column
    rx, rw = MARGIN + Inches(6.66), Inches(5.23)
    tf = textbox(s, rx, y, rw, Inches(0.3))
    para(tf, "2  \u00b7  WHAT IS ONE CALL?   PREDICATES, NOT NAMES",
         size=10.5, color=PURPLE, bold=True, font=FONT_SB, first=True, space_after=0)

    c = card(s, rx, y + Inches(0.34), rw, Inches(1.14), fill=WHITE, line=LINE)
    tf = c.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.24)
    para(tf, "LLM_PREDICATE", size=10, color=PURPLE, bold=True, font=MONO, first=True,
         space_after=5)
    para(tf, "name LIKE 'chat %' OR name = 'chat'\nOR name LIKE '%llm_request%'\nOR name = 'run_sampling_request'".replace("\n", "  "),
         size=10.5, color=INK, font=MONO, space_after=5, line=1.2)
    para(tf, "Matches exactly one span per real model call, per harness \u2014 including only the "
             "outermost of Codex\u2019s five.", size=10.5, color=GRAY, space_after=0, line=1.2)

    c = card(s, rx, y + Inches(1.62), rw, Inches(1.14), fill=WHITE, line=LINE)
    tf = c.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.24)
    para(tf, "TOOL_PREDICATE", size=10, color=PURPLE, bold=True, font=MONO, first=True,
         space_after=5)
    para(tf, "name LIKE 'execute_tool%'  OR  name = 'claude_code.tool'  OR  name = 'handle_tool_call'",
         size=10.5, color=INK, font=MONO, space_after=5, line=1.2)
    para(tf, "Exact match on claude_code.tool \u2014 the .execution and .blocked_on_user children are "
             "deliberately excluded.", size=10.5, color=GRAY, space_after=0, line=1.2)

    c = card(s, rx, y + Inches(2.9), rw, Inches(0.8), fill=BLUE_XL, line=BLUE_XL)
    tf = c.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.24)
    para(tf, "That is what makes any cross-agent number honest. One definition, applied to three "
             "vocabularies \u2014 rather than three definitions producing three numbers.",
         size=11.5, color=NAVY, first=True, space_after=0, line=1.24)

    # bottom
    band = card(s, MARGIN, Inches(5.98), Inches(11.89), Inches(0.7), fill=NAVY, line=NAVY)
    tf = band.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.32)
    para(tf, [("3  \u00b7  The agent\u2019s identity comes from the host\u2019s URI scheme, not its service name.  ",
               dict(bold=True, color=BLUE_L)),
              ("claude:/\u2026 \u2192 Claude, codex:/\u2026 \u2192 Codex, copilotcli:/\u2026 \u2192 Copilot CLI. "
               "The host doesn\u2019t control what resource name an agent picks for itself \u2014 "
               "but it does control the URI it launched it with.",
               dict(color=RGBColor(0xC3, 0xD1, 0xE6)))],
         size=12, first=True, space_after=0, line=1.24)

    footer(s, FOOT + "  \u00b7  appendix")
    notes(s, """
APPENDIX \u2014 not in the running order. Jump here if someone asks how sessions are identified.

Say the pattern first, then show it three times:
"Define the concept once, in SQL, from what the data SHOWS \u2014 never from which harness sent it."

LEFT: session identity is a COALESCE ladder ordered by trust. OTel's standard key first, then
each vendor's own, then a recovered projection for Codex, then trace_id as the last resort.
Point at MAX() and explain it: the conversation id is on the chat span but not on the permission
span or the execute_tool span in the same trace. MAX over the trace means one span knowing the
answer is enough for the whole trace.

RIGHT: 'one model call' and 'one tool call' are predicates, not name lists. The Claude detail is
worth 10 seconds: match the subtree and you treble the count; match only .execution and you
silently drop every permission-denied call \u2014 which are the calls you most want to see.

BOTTOM: this is the agent-host fix from earlier, in one line.
""")
    return s


# ------------------------------------------------------------------ normalization II
def s_norm_tokens(prs):
    s = blank(prs)
    y = header(s, "03  Approach \u00b7 the contribution", "Normalization \u2014 the token problem",
               "One worked example of the pattern: three incompatible counting schemes, "
               "one honest number.")

    # three convention cards
    cw = Inches(3.83)
    gap = Inches(0.2)
    convs = [
        ("OTel semconv", "Copilot \u00b7 Codex", BLUE,
         "cache_read \u2286 input",
         "prompt = input", "already counted"),
        ("Anthropic", "Claude Code", AMBER,
         "cache is ADDITIVE, not a subset",
         "prompt = input + cache_read", "         + cache_creation"),
        ("Rollup spans", "Copilot \u00b7 Codex", MAGENTA,
         "parents re-report their children",
         "one call counted twice", "silently"),
    ]
    for i, (title, who, col, rule, l1, l2) in enumerate(convs):
        x = MARGIN + (cw + gap) * i
        c = card(s, x, y, cw, Inches(1.72), fill=WHITE, line=LINE)
        rect(s, x, y, cw, Inches(0.055), fill=col)
        tf = c.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = tf.margin_right = Inches(0.24)
        tf.margin_top = Inches(0.2)
        para(tf, [(title, dict(size=13.5, color=NAVY, bold=True)),
                  ("    " + who, dict(size=10, color=col, bold=True))],
             first=True, space_after=6)
        para(tf, rule, size=11.5, color=col, bold=True, space_after=7, line=1.2)
        para(tf, l1, size=10.5, color=INK, font=MONO, space_after=4, line=1.2)
        para(tf, l2, size=10.5, color=GRAY, font=MONO, space_after=0, line=1.2)

    # resolution
    tf = textbox(s, MARGIN, y + Inches(1.96), Inches(11.89), Inches(0.3))
    para(tf, "HOW IT RESOLVES \u2014 ONE FACT TABLE, WRITTEN AT INSERT TIME", size=10.5,
         color=TEAL, bold=True, font=FONT_SB, first=True, space_after=0)

    grey = RGBColor(0x7A, 0x93, 0xB8)
    code_block(s, MARGIN, y + Inches(2.3), Inches(6.4), Inches(1.7), [
        [("token_facts", dict(color=BLUE_L, bold=True)),
         ("  \u2014 one row per span, exempt from retention", dict(color=grey))],
        [("  input \u00b7 output \u00b7 cache_read \u00b7 cache_creation \u00b7 model", {})],
        [("  is_additive", dict(color=AMBER)), ("   \u2190 set when Anthropic\u2019s bare", dict(color=grey))],
        [("                    cache_*_tokens keys appear", dict(color=grey))],
        [("", {})],
        [("CASE WHEN", dict(color=BLUE_L)), (" is_additive ", {}), ("THEN", dict(color=BLUE_L)),
         (" input+read+creation", {})],
        [("                        ", {}), ("ELSE", dict(color=BLUE_L)), (" input ", {}),
         ("END", dict(color=BLUE_L)), (" AS prompt_tokens", {})],
    ], size=11)

    rx = MARGIN + Inches(6.66)
    items = [
        ("One key cascade", "Six vendor spellings per quantity collapse into one number.", TEAL),
        ("Rollups excluded", "Only real model calls contribute \u2014 parent spans never do.", MAGENTA),
        ("Host spans contribute nothing", "Correlation wrappers can\u2019t inflate a total.", PURPLE),
        ("Missing model? Walk up.", "Codex puts tokens below the span that names the model.", BLUE),
    ]
    yy = y + Inches(1.9)
    for title, body, col in items:
        c = card(s, rx, yy, Inches(5.23), Inches(0.68), fill=WHITE, line=LINE)
        rect(s, rx, yy, Inches(0.055), Inches(0.68), fill=col)
        tf = c.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.24)
        tf.margin_right = Inches(0.18)
        para(tf, title, size=11.5, color=NAVY, bold=True, first=True, space_after=2)
        para(tf, body, size=9.5, color=GRAY, space_after=0, line=1.16)
        yy += Inches(0.76)

    band = card(s, MARGIN, Inches(6.06), Inches(6.4), Inches(0.6), fill=BLUE_XL, line=BLUE_XL)
    tf = band.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.26)
    para(tf, [("One test asserts all three at once: ", dict(color=NAVY)),
              ("292 input \u00b7 190 cached \u00b7 35 output \u00b7 3 calls", dict(color=BLUE, bold=True, font=MONO))],
         size=11.5, first=True, space_after=0)

    footer(s, FOOT)
    notes(s, """
SCRIPT

This is the only normalization slide in the running order \u2014 it carries the whole "this is the
contribution" argument. Give it a full beat. Identity and transcripts are in the appendix.

"I'll show you one of these problems properly, because they all have the same shape. Tokens \u2014
because everyone thinks token counting is trivial."

[the three cards]
"There are three incompatible counting schemes running on one developer's machine at the same
time."

"Copilot and Codex follow the OpenTelemetry convention: cache reads are a subset of your input
tokens. Already counted. Don't add them."

"Anthropic does the opposite. Cache reads are additive \u2014 they sit on top of input. Same concept,
opposite arithmetic. If you get that backwards, you under-report Claude's prompt tokens by three
or four times."

"And the third one is sneaky. Some parent spans re-report everything their children already
reported. Sum those naively and a single model call gets counted twice."

[the resolution]
"The fix is one fact table, written at insert time. Every span gets a row, and a flag saying
whether this vendor's cache numbers are additive. Then the query expands conditionally."

"The important detail is how that flag gets set: by the SHAPE of the attributes \u2014 the presence
of Anthropic's bare cache key names \u2014 never by 'if the vendor is Claude'. That's the difference
between code that survives the next release and code that doesn't."

[the test line]
"And one test asserts all three at once from a single fixture. If any vendor's convention drifts,
that test fails immediately."

[close on the generalisation \u2014 this is what makes one example stand for the layer]
"Identity, tool counts and transcripts all needed exactly this treatment. Same pattern every
time: define the concept once, in SQL, from what the data shows rather than from who sent it.
And this is the shape that falls out the other end."

\u2500\u2500\u2500 DETAIL FOR Q&A \u2500\u2500\u2500

\u2022 The cache HIT RATE has a different denominator too \u2014 read/input under OTel, but
  read/(input+read+creation) under Anthropic. Same trap, second order.
\u2022 The rollup spans specifically: Copilot's invoke_agent aggregates every subagent, and Codex has
  a turn-level span that aggregates the whole turn. Both are excluded by only counting spans whose
  operation name marks them as a real chat completion.
\u2022 The key cascade collapses up to six vendor spellings per quantity, and it's exported from the
  shared types package so the store and the engine physically cannot disagree.
\u2022 "Walk up": Codex puts token counts on a span three levels below the one naming the model, so a
  recursive query climbs the ancestor chain to attribute them.
\u2022 token_facts is exempt from retention pruning \u2014 it outlives the raw spans, which is what makes
  nine days of trend data possible without storing nine days of spans.
""")
    return s


# ------------------------------------------------------------------ normalization III
def s_norm_transcripts(prs):
    s = blank(prs)
    y = header(s, "Appendix \u00b7 normalization",
               "Transcripts \u2014 three sources, one shape",
               "The renderer and the twelve LM tools have no idea which harness they are "
               "looking at. That is the whole idea.")

    src = [("Span attributes", "Copilot / Agent Host", BLUE,
            "gen_ai.output.messages  \u00b7  gen_ai.input.messages\nrichest source; tool_call parts already structured"),
           ("claudeLogTurns()", "Claude Code", AMBER,
            "rebuilds turns from the log stream, delimited by\napi_request records \u2014 logged when the round trip finishes,\nso a record can arrive before the call it closes"),
           ("codexLogTurns()", "Codex", MAGENTA,
            "prompts and tool results from log records; token counts\njoined by conversation.id, because response.completed\ncarries no trace id at all")]
    cw = Inches(3.83)
    for i, (title, who, col, body) in enumerate(src):
        x = MARGIN + (cw + Inches(0.2)) * i
        c = card(s, x, y, cw, Inches(1.5), fill=WHITE, line=col, radius=0.05)
        c.line.width = Pt(1.5)
        tf = c.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = tf.margin_right = Inches(0.22)
        tf.margin_top = Inches(0.18)
        para(tf, [(title, dict(size=12.5, color=NAVY, bold=True, font=MONO)),
                  ("   " + who, dict(size=10, color=col, bold=True))],
             first=True, space_after=6)
        para(tf, body.replace("\n", " "), size=10.5, color=GRAY, space_after=0, line=1.22)
        down_arrow(s, x + cw / 2, y + Inches(1.58), color=col, w=Inches(0.18), h=Inches(0.26))

    # converge
    sh = card(s, MARGIN, y + Inches(1.98), Inches(11.89), Inches(0.82), fill=NAVY, line=NAVY)
    tf = sh.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.36)
    para(tf, [("SessionMessageTurn", dict(size=16, color=WHITE, bold=True, font=MONO)),
              ("     traceId \u00b7 spanId \u00b7 sourceSpanId \u00b7 model \u00b7 outputMessages \u00b7 inputPreview \u00b7 "
               "systemInstructions \u00b7 details[] \u00b7 isSubagent",
               dict(size=10.5, color=RGBColor(0x9E, 0xB4, 0xD6), font=MONO))],
         first=True, space_after=3)
    para(tf, "One canonical shape. No renderer and no LM tool contains a single vendor branch.",
         size=11.5, color=BLUE_L, space_after=0)

    # three more rules
    rules = [
        ("Real sessions vs. artefacts", MAGENTA,
         "BACKGROUND_TRACE_FILTER keeps a session only if it shows agent activity, a captured "
         "user prompt, or a host-assigned title \u2014 261 phantom sessions become the 6 real ones."),
        ("Written from shape, not from vendor", GREEN,
         "ECHO_TRACE is defined as \u201chas content, no prompt, no round trip, no model-call span\u201d "
         "\u2014 never as \u201cif it\u2019s Codex\u201d. That is why it still holds after three upstream releases."),
        ("Prompts are cleaned before they\u2019re shown", TEAL,
         "<system-reminder> blocks, injected repository context and Copilot\u2019s <userRequest> "
         "envelope are stripped. A 120-char label whose first 66 chars are an injected "
         "<current_datetime> stamp is not a label."),
    ]
    yy = y + Inches(2.96)
    for i, (title, col, body) in enumerate(rules):
        x = MARGIN + (cw + Inches(0.2)) * i
        c = card(s, x, yy, cw, Inches(1.6), fill=WHITE, line=LINE)
        rect(s, x, yy, cw, Inches(0.055), fill=col)
        tf = c.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.22)
        para(tf, title, size=12, color=NAVY, bold=True, first=True, space_after=5, line=1.15)
        para(tf, body, size=10.5, color=GRAY, space_after=0, line=1.22)

    footer(s, FOOT + "  \u00b7  appendix")
    notes(s, """
APPENDIX \u2014 not in the running order. Jump here if someone asks how the transcript is rebuilt,
or why Codex transcripts are missing the model's replies.

Top row: three completely different acquisition strategies. Copilot hands you structured
messages on a span attribute. Claude requires rebuilding turns from a log stream. Codex requires
rebuilding them from logs AND joining tokens by conversation id, because its response.completed
records carry no trace id at all \u2014 joining by trace matched nothing, and for a while no Codex
transcript reported a single token.

The Claude detail is my favourite: api_request is logged when the round trip FINISHES. So a
tool_decision log can legitimately arrive before the api_request that owns it. The builder
adopts an already-open unclaimed draft instead of minting a new call \u2014 otherwise a seven-call
agent loop collapses into one turn with every tool piled onto the last call.

Then the navy bar, and this is the sentence:
"All three reshape into the same SessionMessageTurn, so no renderer and no LM tool knows which
harness it is looking at."

Bottom-middle card is the design principle and it becomes Learning #1: written from shape, not
from vendor. Harness-specific code is a dated cheque.
""")
    return s


# ------------------------------------------------------------------ the session model
def s_session_model(prs):
    s = blank(prs)
    y = header(s, "03  Approach \u00b7 the output", "One model out \u2014 what a session actually is",
               "The shape everything above produces, and the only shape the panel "
               "and the tools ever see.")

    # ---------------- left: nested model
    lx, lw = MARGIN, Inches(6.9)
    outer = card(s, lx, y + Inches(0.04), lw, Inches(3.98), fill=WASH, line=LINE)
    strip = card(s, lx, y + Inches(0.04), lw, Inches(0.46), fill=NAVY, line=NAVY)
    tf = strip.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.26)
    para(tf, [("SESSION", dict(color=WHITE, bold=True, size=11.5, font=FONT_SB)),
              ("     agent \u00b7 outcome \u00b7 duration \u00b7 total tokens \u00b7 cache hit rate",
               dict(color=RGBColor(0x93, 0xA6, 0xC4), size=10.5))],
         first=True, space_after=0)

    turns = [("TURN 1", "\u201cfix the flaky test\u201d",
              [("model call", PURPLE, 2.05), ("tool call \u00b7 read_file", AMBER, 1.5),
               ("tool call \u00b7 edit_file", AMBER, 1.5)]),
             ("TURN 2", "\u201cnow add a regression test\u201d",
              [("model call", PURPLE, 2.05), ("tool call \u00b7 run_tests", AMBER, 1.7)])]
    ty = y + Inches(0.66)
    for label, quote, pills in turns:
        tb = card(s, lx + Inches(0.24), ty, lw - Inches(0.48), Inches(1.5),
                  fill=WHITE, line=LINE)
        rect(s, lx + Inches(0.24), ty, lw - Inches(0.48), Inches(0.05), fill=BLUE)
        tf = textbox(s, lx + Inches(0.44), ty + Inches(0.16), lw - Inches(0.9), Inches(0.32))
        para(tf, [(label, dict(color=BLUE, bold=True, size=10.5, font=FONT_SB)),
                  ("     " + quote, dict(color=GRAY, size=10.5, italic=True))],
             first=True, space_after=0)
        px = lx + Inches(0.44)
        for name, col, pw in pills:
            p = card(s, px, ty + Inches(0.66), Inches(pw), Inches(0.5),
                     fill=WHITE, line=col)
            rect(s, px, ty + Inches(0.66), Inches(0.05), Inches(0.5), fill=col)
            tf = p.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Inches(0.18)
            tf.margin_right = Inches(0.1)
            para(tf, name, size=9.5, color=NAVY, bold=True, first=True, space_after=0,
                 line=1.1)
            px += Inches(pw) + Inches(0.14)
        ty += Inches(1.66)

    tf = textbox(s, lx, y + Inches(4.14), lw, Inches(0.4))
    para(tf, [("Every number in the panel is an aggregate of this tree. ",
               dict(color=BLUE, bold=True)),
              ("Nothing is stored twice.", dict(color=GRAY))],
         size=11, first=True, space_after=0)

    # ---------------- right: one concept, three spellings
    rx = MARGIN + lw + Inches(0.23)
    rw = Inches(4.76)
    tf = textbox(s, rx, y + Inches(0.04), rw, Inches(0.3))
    para(tf, "ONE CONCEPT, THREE SPELLINGS", size=10.5, color=MAGENTA, bold=True,
         font=FONT_SB, first=True, space_after=0)

    rows = [("Session", "conversation.id  \u00b7  session.id  \u00b7  chat_session_id  \u00b7  or nothing"),
            ("Turn", "one chat span  \u00b7  api_request  \u00b7  run_sampling_request"),
            ("Model call", "one span  \u00b7  one span  \u00b7  outermost of five"),
            ("Tool call", "execute_tool%  \u00b7  claude_code.tool  \u00b7  handle_tool_call"),
            ("Tokens", "on the span  \u00b7  additive cache keys  \u00b7  walk up the tree")]
    ry = y + Inches(0.42)
    for name, srcs in rows:
        c = card(s, rx, ry, rw, Inches(0.72), fill=WHITE, line=LINE)
        rect(s, rx, ry, Inches(0.055), Inches(0.72), fill=TEAL)
        tf = textbox(s, rx + Inches(0.24), ry + Inches(0.1), rw - Inches(0.42), Inches(0.62))
        para(tf, name, size=11.5, color=NAVY, bold=True, first=True, space_after=2)
        para(tf, srcs, size=8.5, color=GRAY, font=MONO, space_after=0, line=1.16)
        ry += Inches(0.78)

    band = card(s, rx, ry + Inches(0.04), rw, Inches(0.54), fill=NAVY, line=NAVY)
    tf = band.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.24)
    para(tf, "Left side is the product. Right side is why it took a summer.",
         size=11, color=WHITE, bold=True, first=True, space_after=0, line=1.15)

    footer(s, FOOT)
    notes(s, """
SCRIPT

Thirty seconds. This is the payoff slide for the whole Approach section \u2014 it's what all the
normalization is FOR. Do not skip it, even if you're behind.

"So this is what comes out the other end."

[point at the left]
"A session contains turns. A turn is one thing you asked for, plus everything the agent did to
answer it \u2014 the model calls and the tool calls. Tokens roll up."

"That's it. That's the whole model."

"And every number you'll see in the demo is an aggregate of this tree. Nothing is stored twice."

[sweep the right column once \u2014 do NOT read it]
"And that's the same comparison table I opened this section with, collapsed down into five
concepts."

[the navy band \u2014 land it and stop]
"The left side is the product. The right side is why it took a summer."

\u2500\u2500\u2500 DETAIL FOR Q&A \u2500\u2500\u2500

\u2022 "What exactly is a turn?" \u2014 one user request plus every model and tool call it triggered,
  bounded by the next user request. That definition is the thing the three harnesses do not
  agree on, and the thing everything downstream depends on.
\u2022 "Nothing is stored twice" matters because it means adding a new view is a query, not a
  migration. The session, the timeline and the transcript are three reads of the same rows.
\u2022 If asked how sessions are identified: that's an appendix slide \u2014 a trust-ordered ladder that
  tries the standard key first, then each vendor's own, then a recovered projection for Codex.
""")
    return s